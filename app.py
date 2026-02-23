# app.py - GEO Pulse: Bloomberg terminal for the GEO/AEO category
import io
import json
import math
import os
import re
from datetime import datetime, timedelta
from collections import defaultdict, Counter

import streamlit as st
import pandas as pd
import altair as alt
import plotly.graph_objects as go
from docx import Document as DocxDocument
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pipeline.enrich import run_enrichment

st.set_page_config(page_title="GEO Pulse", page_icon="📡", layout="wide")

# ---------------------------------------------------------------------------
# Gist brand theme
# ---------------------------------------------------------------------------

st.markdown("""<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Sans:wght@400;500;700&family=DM+Mono:wght@400;500&display=swap');

/* Primary button overrides */
.stButton > button[kind="primary"],
.stFormSubmitButton > button,
.stFormSubmitButton > button[kind="primary"],
.stDownloadButton > button {
  background-color: #0E3B7E !important;
  color: #F8F4EB !important;
  font-family: 'DM Sans', sans-serif !important;
  font-weight: 500;
  border: none !important;
  border-radius: 0px !important;
}
.stButton > button[kind="primary"]:hover,
.stFormSubmitButton > button:hover,
.stDownloadButton > button:hover {
  background-color: #FF9D1C !important;
  color: #0A0A0A !important;
}

/* Secondary buttons — filled navy by default */
.stButton > button[kind="secondary"] {
  background-color: #0E3B7E !important;
  color: #F8F4EB !important;
  font-family: 'DM Sans', sans-serif !important;
  font-weight: 500;
  border: none !important;
  border-radius: 0px !important;
}
.stButton > button[kind="secondary"]:hover {
  background-color: #FF9D1C !important;
  color: #0A0A0A !important;
}


html, body, [class*="css"] {
  font-family: 'DM Sans', sans-serif;
  background-color: #F8F4EB;
  color: #0A0A0A;
}
.stApp { background-color: #F8F4EB; }

h1 { font-family: 'DM Sans', sans-serif; font-weight: 700; color: #0A0A0A; }
h2, h3 { font-family: 'DM Sans', sans-serif; font-weight: 500; color: #0A0A0A; }

.stMetric label {
  font-family: 'DM Mono', monospace;
  font-size: 11px; text-transform: uppercase;
  letter-spacing: 0.08em; color: #0A0A0A;
}
.stMetric [data-testid="metric-value"] {
  font-family: 'DM Sans', sans-serif; font-weight: 700; color: #0A0A0A;
}

.stTabs [data-baseweb="tab"] {
  font-family: 'DM Mono', monospace; font-size: 12px;
  text-transform: uppercase; letter-spacing: 0.06em; color: #0A0A0A;
}
.stTabs [aria-selected="true"] { border-bottom: 2px solid #0E3B7E; color: #0A0A0A; }

.streamlit-expander {
  border: 1px solid #D1CFBA; border-radius: 0px; background-color: #F8F4EB;
}
.streamlit-expander [data-testid="stExpanderDetails"] {
  background-color: #FFFFFF; padding: 16px; border-top: 1px solid #D1CFBA;
}

.element-container div[data-testid="stMarkdownContainer"] {
  font-family: 'DM Sans', sans-serif;
}

code {
  font-family: 'DM Mono', monospace;
  background-color: #D1CFBA; color: #0A0A0A;
  padding: 2px 6px; border-radius: 0px; font-size: 11px;
}

[data-testid="stAlert"] { border-radius: 0px; border-left: 4px solid #FF9D1C; }

.css-1d391kg { background-color: #D1CFBA; }

.stSelectbox > div > div {
  border-radius: 0px; border: 1px solid #D1CFBA;
  font-family: 'DM Sans', sans-serif;
}

.stMultiSelect span[data-baseweb="tag"] {
  background-color: #0E3B7E; color: #F8F4EB;
  border-radius: 0px; font-family: 'DM Mono', monospace; font-size: 11px;
}

[data-testid="stSidebar"] { display: none !important; }
[data-testid="stSidebarCollapsedControl"] { display: none !important; }
</style>""", unsafe_allow_html=True)

# ---------------------------------------------------------------------------
# ---------------------------------------------------------------------------
# Data loading
# ---------------------------------------------------------------------------

DATA_DIR = "data"
INSIGHTS_PATH = os.path.join(DATA_DIR, "enriched_insights.json")
TRENDS_PATH = os.path.join(DATA_DIR, "trends.json")
SOURCES_PATH = os.path.join(DATA_DIR, "discovered_sources.json")
CLUSTERS_PATH = os.path.join(DATA_DIR, "clusters.json")
COMPANIES_PATH = "config/companies.json"
RUN_LOG_PATH = os.path.join(DATA_DIR, "run_log.json")


@st.cache_data(ttl=300)
def load_insights():
    if os.path.exists(INSIGHTS_PATH):
        with open(INSIGHTS_PATH, "r") as f:
            return json.load(f)
    return []


@st.cache_data(ttl=300)
def load_trends():
    if os.path.exists(TRENDS_PATH):
        with open(TRENDS_PATH, "r") as f:
            return json.load(f)
    return {}


@st.cache_data(ttl=300)
def load_discovered_sources():
    if os.path.exists(SOURCES_PATH):
        with open(SOURCES_PATH, "r") as f:
            return json.load(f)
    return []


@st.cache_data(ttl=600)
def load_companies():
    if os.path.exists(COMPANIES_PATH):
        with open(COMPANIES_PATH, "r") as f:
            return json.load(f)
    return {"own_brands": [], "competitors": []}


def load_run_log():
    if os.path.exists(RUN_LOG_PATH):
        with open(RUN_LOG_PATH, "r") as f:
            return json.load(f)
    return []


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

SOURCE_LIST = "Reddit, Hacker News, G2, Product Hunt, Google News, Search Engine Journal, Search Engine Land, Digiday, AdExchanger"


def _source_badge(source):
    """Plain text source label for pill display."""
    for key in ("Reddit", "Hacker News", "Slack", "Product Hunt", "G2", "News", "RSS"):
        if key.lower() in source.lower():
            return key
    return source.strip() or "Source"


_TITLE_BLOCKLIST = re.compile(
    r"^\s*\[(dead|flagged|deleted)\]\s*$"
    r"|who is hiring"
    r"|who.s hiring"
    r"|ask hn:.*hiring"
    r"|hiring thread"
    r"|freelancer.*seeking"
    r"|monthly.*job"
    r"|why is the sky"
    r"|50 years ago"
    r"|bill gates.*birthday"
    r"|what color is"
    r"|how to lose weight"
    r"|best recipe for"
    r"|weather forecast"
    r"|horoscope"
    r"|lottery results"
    r"|sports score"
    r"|celebrity gossip",
    re.I,
)

_OFFTOPIC_BLOCKLIST = re.compile(
    r"cooking|recipe|fitness|workout|horoscope|astrology"
    r"|weight loss|diet plan|celebrity|gossip|movie review"
    r"|sports score|game recap|lottery|weather forecast"
    r"|real estate listing|mortgage rate|car review"
    r"|travel deal|vacation package|dating advice",
    re.I,
)

_MIN_RELEVANCE_SCORE = 1.5


def _is_displayable_post(insight):
    """Filter out dead links, job posts, and empty titles."""
    title = (insight.get("title") or "").strip()
    if not title:
        return False
    if _TITLE_BLOCKLIST.search(title):
        return False
    return True


def _dedup_insights(posts):
    """Remove duplicates by URL, keeping the most enriched version."""
    by_url = {}
    no_url = []
    for p in posts:
        url = (p.get("url") or "").strip()
        if not url:
            no_url.append(p)
            continue
        existing = by_url.get(url)
        if existing is None:
            by_url[url] = p
        else:
            new_score = len(p.get("companies_mentioned", [])) + len(p.get("entity_tags", []))
            old_score = len(existing.get("companies_mentioned", [])) + len(existing.get("entity_tags", []))
            if new_score > old_score:
                by_url[url] = p
    return list(by_url.values()) + no_url


def _time_ago(date_str):
    try:
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        days = (datetime.now() - dt).days
        if days == 0:
            return "today"
        if days == 1:
            return "1d ago"
        if days < 7:
            return f"{days}d ago"
        if days < 30:
            return f"{days // 7}w ago"
        return f"{days // 30}mo ago"
    except (ValueError, TypeError):
        return ""


def _relevance_score(insight):
    """Composite score: source quality x social signal x recency x enrichment signals."""
    quality = insight.get("source_quality", 2)
    score = insight.get("score", 0)
    comments = insight.get("num_comments", 0)
    social = 1 + math.log1p(min(score, 500) + comments)

    try:
        dt = datetime.strptime(insight.get("post_date", ""), "%Y-%m-%d")
        days_old = max(0, (datetime.now() - dt).days)
        recency = max(0.1, 1 - days_old / 60)
    except (ValueError, TypeError):
        recency = 0.3

    # Boost for enrichment signals
    enrichment = 1.0
    if insight.get("companies_mentioned"):
        enrichment += 0.5 * len(insight["companies_mentioned"])
    if insight.get("entity_tags"):
        enrichment += 0.3 * len(insight["entity_tags"])
    if insight.get("is_competitive_intel"):
        enrichment += 2.0
    if insight.get("is_feature_request"):
        enrichment += 1.0
    if insight.get("is_buyer_voice"):
        enrichment += 0.5

    return quality * social * recency * enrichment


def _relevance_sentence(insight, for_company=None):
    """One-sentence brief explaining WHY this signal matters for GEO/AEO practitioners.

    When for_company is set, that company is placed first in rationale text
    so the sentence makes sense on that company's card.
    """
    # If the article has a summary/lede, prefer that
    summary = (insight.get("summary") or insight.get("lede") or "").strip()
    if summary and len(summary) > 20:
        return summary[:200]

    tags = insight.get("entity_tags", [])
    companies = list(insight.get("companies_mentioned", []))
    features = insight.get("features_mentioned", [])
    sentiment = insight.get("sentiment", "neutral")
    source = insight.get("source", "")
    feat_ctx = f" in {', '.join(features[:2])}" if features else ""

    # Reorder so for_company appears first
    if for_company and for_company in companies:
        companies = [for_company] + [c for c in companies if c != for_company]

    comp_str = ", ".join(companies[:3]) if companies else ""

    if insight.get("is_competitive_intel"):
        if len(companies) >= 2:
            return f"Practitioners are comparing {companies[0]} and {', '.join(companies[1:3])}{feat_ctx}, signaling active buyer evaluation in this space."
        return f"Competitive intelligence signal{feat_ctx} worth tracking for positioning decisions."
    if insight.get("is_feature_request"):
        voice = "buyers" if insight.get("is_buyer_voice") else "users"
        if features:
            return f"Market demand signal: {voice} are asking for {', '.join(features[:2])}, indicating an unmet need in the category."
        return f"Market demand signal: {voice} are requesting new capabilities, revealing gaps in current tooling."
    if "funding_news" in tags:
        target = f"{companies[0]} " if companies else "a GEO/AEO player "
        return f"Capital is flowing: {target}raised funding, signaling investor confidence in this category."
    if "product_launch" in tags:
        target = f"{companies[0]} " if companies else "A competitor "
        return f"{target}launched new capabilities{feat_ctx}. Watch for positioning shifts and buyer reaction."
    if "complaint" in tags and companies:
        return f"Users are flagging pain points with {companies[0]}{feat_ctx}. Potential differentiation opportunity."
    if "praise" in tags and companies:
        return f"{companies[0]} is earning positive sentiment{feat_ctx}. Worth studying what they are doing right."
    if insight.get("is_buyer_voice"):
        if companies:
            return f"Active buyer evaluating {comp_str}{feat_ctx}. Direct demand signal for the category."
        return f"Buyer actively evaluating GEO/AEO tools{feat_ctx}. Direct demand signal."
    if insight.get("is_founder_voice"):
        target = f" at {companies[0]}" if companies else ""
        return f"Founder perspective{target}{feat_ctx}. Insider signal on product direction."
    if insight.get("is_analyst_voice"):
        return f"Analyst or trade press coverage{feat_ctx}. Shapes buyer perception and category narrative."
    if companies and features:
        sent = f" Sentiment trending {sentiment}." if sentiment != "neutral" else ""
        return f"Market discussion about {', '.join(companies[:2])} in context of {', '.join(features[:2])}.{sent}"
    if companies:
        sent = f" Sentiment trending {sentiment}." if sentiment != "neutral" else ""
        return f"Industry conversation involving {', '.join(companies[:2])} in the GEO/AEO space.{sent}"
    return ""


def _keywords_for_card(insight):
    """Extract 3-5 display keywords from enrichment data."""
    kws = []
    for tag in insight.get("entity_tags", []):
        kws.append(tag.replace("_", " "))
    for feat in insight.get("features_mentioned", []):
        if feat not in kws:
            kws.append(feat)
    if insight.get("is_buyer_voice") and "buyer voice" not in kws:
        kws.append("buyer voice")
    if insight.get("is_founder_voice") and "founder voice" not in kws:
        kws.append("founder voice")
    return kws[:5]


@st.cache_data(ttl=300)
def _get_new_companies(_insights_json):
    """Pre-compute which companies were first seen in last 7 days."""
    insights_list = json.loads(_insights_json)
    company_oldest = {}
    now = datetime.now()
    for i in insights_list:
        date_str = i.get("post_date", "")
        for comp in i.get("companies_mentioned", []):
            if comp not in company_oldest or (date_str and date_str < company_oldest.get(comp, "")):
                company_oldest[comp] = date_str
    new_set = set()
    for comp, oldest in company_oldest.items():
        try:
            dt = datetime.strptime(oldest, "%Y-%m-%d")
            if (now - dt).days <= 7:
                new_set.add(comp)
        except (ValueError, TypeError):
            pass
    return new_set


CATEGORY_LABELS = {
    "geo_measurement": "GEO/AEO Measurement",
    "geo_as_a_service_publisher": "GEO-as-a-Service / Publisher",
    "ai_attribution": "AI Attribution",
    "seo_suite": "SEO Suite",
    "content_optimization": "Content Optimization",
    "ai_content": "AI Content",
    "brand_monitoring": "Brand Monitoring",
    "competitive_intel": "Competitive Intel",
    "sales_intel": "Sales Intelligence",
}


def _company_positioning(comp_data):
    """One-line positioning from company metadata."""
    cat = comp_data.get("category", "")
    notes = comp_data.get("notes", "")
    if notes:
        return notes
    return CATEGORY_LABELS.get(cat, cat.replace("_", " ").title())


# ---------------------------------------------------------------------------
# Share link helper
# ---------------------------------------------------------------------------

def _share_button(section_id, label="Share", key_suffix="", extra_params=None):
    """Render a share button that copies a deep link to clipboard."""
    btn_key = f"share_{section_id}_{key_suffix}" if key_suffix else f"share_{section_id}"
    if st.button(f"🔗 {label}", key=btn_key, type="secondary"):
        params = {"section": section_id}
        if extra_params:
            params.update(extra_params)
        qs = "&".join(f"{k}={v}" for k, v in params.items())
        import streamlit.components.v1 as components
        full_url = f"?{qs}"
        components.html(
            f"""<script>
            navigator.clipboard.writeText(window.location.origin + window.location.pathname + '{full_url}');
            </script>
            <div style="background:#0A0A0A;color:#F8F4EB;padding:6px 12px;
            font-size:0.8rem;text-align:center;">Link copied to clipboard</div>""",
            height=36,
        )


# ---------------------------------------------------------------------------
# Citation card helper
# ---------------------------------------------------------------------------



# ---------------------------------------------------------------------------
# Export generators
# ---------------------------------------------------------------------------

_METHODOLOGY_NOTE = (
    "Data collected via automated pipeline monitoring 11 public sources. "
    "Signals filtered for GEO/AEO relevance. Sentiment scored by LLM enrichment. "
    "Competitor coverage based on public mentions only. Absence of data does not "
    "indicate absence of a feature. Signal counts reflect volume of public conversation, "
    "not product capability assessments."
)

_CONFIDENCE_FOOTNOTE = (
    "Confidence score reflects source diversity and signal volume, "
    "not predictive accuracy."
)

_SOURCE_DESCRIPTIONS = [
    ("Hacker News", "Founder and builder discussions"),
    ("Reddit", "Practitioner community signals"),
    ("G2", "Buyer reviews and comparisons"),
    ("Product Hunt", "Product launches and feature announcements"),
    ("Google News", "Industry news aggregation"),
    ("Search Engine Journal", "SEO/GEO trade analysis"),
    ("Search Engine Land", "Search marketing news"),
    ("Digiday", "Digital media and marketing coverage"),
    ("AdExchanger", "Ad-tech and martech industry reporting"),
    ("RSS feeds", "Curated trade and blog sources"),
    ("Slack communities", "Private practitioner channels (opted-in)"),
]


def _docx_source_caption(doc, total_signals, date_str=None):
    """Add a small gray italic data-source line to a docx document."""
    if not date_str:
        date_str = datetime.now().strftime("%Y-%m-%d")
    p = doc.add_paragraph()
    run = p.add_run(
        f"Source: GEO Pulse | {total_signals:,} signals from Hacker News, Reddit, G2, "
        f"Product Hunt, and trade press | Data as of {date_str}"
    )
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.italic = True
    return p


def _docx_confidence_footnote(doc):
    """Add a confidence disclosure footnote to a docx document."""
    p = doc.add_paragraph()
    run = p.add_run(f"* {_CONFIDENCE_FOOTNOTE}")
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.italic = True
    return p


def _docx_methodology_appendix(doc, total_signals):
    """Add a full methodology section to a docx document."""
    doc.add_heading("Methodology", level=1)
    doc.add_paragraph(_METHODOLOGY_NOTE)

    doc.add_heading("Sources Monitored", level=2)
    for name, desc in _SOURCE_DESCRIPTIONS:
        p = doc.add_paragraph(style="List Bullet")
        run_name = p.add_run(f"{name}")
        run_name.bold = True
        p.add_run(f", {desc}")

    doc.add_paragraph("")
    p = doc.add_paragraph()
    run = p.add_run(
        f"Pipeline refreshes every 6 hours. Total signals in database: {total_signals:,}."
    )
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _pptx_slide_footer(slide, date_str, total_signals):
    """Add bottom-left corner annotation to a pptx slide."""
    box = slide.shapes.add_textbox(
        PptxInches(0.3), PptxInches(7.0), PptxInches(6), PptxInches(0.4)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"GEO Pulse | {date_str} | {total_signals:,} signals analyzed"
    p.font.size = PptxPt(9)
    p.font.color.rgb = PptxRGBColor(0x99, 0x99, 0x99)


def _pptx_callout_box(slide, bullets, top_inches=5.6):
    """Add a text box with bullet callouts below a chart image."""
    box = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(top_inches), PptxInches(12), PptxInches(1.5)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = PptxPt(12)
        p.font.color.rgb = PptxRGBColor(0x33, 0x33, 0x33)


def _docx_add_hyperlink(paragraph, text, url, font_size=None, bold=False, color="336699"):
    """Add a clickable hyperlink to a python-docx paragraph."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    part = paragraph.part
    r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    c = OxmlElement("w:color")
    c.set(qn("w:val"), color)
    rPr.append(c)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    if bold:
        rPr.append(OxmlElement("w:b"))
    if font_size:
        sz = OxmlElement("w:sz")
        sz.set(qn("w:val"), str(int(font_size * 2)))
        rPr.append(sz)
        szCs = OxmlElement("w:szCs")
        szCs.set(qn("w:val"), str(int(font_size * 2)))
        rPr.append(szCs)
    new_run.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    t.set(qn("xml:space"), "preserve")
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


# Signal of the Week data (mirrors the hardcoded editorial pick in the Live Feed tab)
_SOTW = {
    "company": "Future",
    "title": "Future PLC launches GEO-as-a-Service division",
    "url": "https://pressgazette.co.uk/marketing/future-leveragess-high-visibility-on-chatgpt-by-offering-geo-as-a-service/",
    "date": "2026-02-20",
    "source": "Press Gazette",
    "brief": (
        "Future PLC, publisher of TechRadar and Tom's Guide and the most-cited "
        "publisher domain on ChatGPT globally, has launched a commercial GEO "
        "optimization division selling AI visibility campaigns to brand clients. "
        "They delivered a 33% ChatGPT visibility uplift for Samsung and hold a "
        "direct content deal with OpenAI. This is the first major media publisher "
        "to productize GEO expertise, signaling the category is moving mainstream."
    ),
}


def _export_research_report(insights, company_meta, opportunity_data, selected_comps, comp_stats):
    """Generate a Research Report .docx styled as an analyst newsletter."""
    doc = DocxDocument()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    total_signals = len(insights)
    date_str = datetime.now().strftime("%Y-%m-%d")
    date_display = datetime.now().strftime("%B %d, %Y")
    _now = datetime.now()
    _7d_ago = (_now - timedelta(days=7)).strftime("%Y-%m-%d")
    _14d_ago = (_now - timedelta(days=14)).strftime("%Y-%m-%d")
    _30d_ago = (_now - timedelta(days=30)).strftime("%Y-%m-%d")
    _90d_ago = (_now - timedelta(days=90)).strftime("%Y-%m-%d")

    # ── Cover Page ──────────────────────────────────────────────────────────
    title_h = doc.add_heading("GEO Pulse Weekly Intelligence", level=0)
    title_h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = sub.add_run(date_display)
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x64, 0x64, 0x64)
    tag = doc.add_paragraph()
    tag.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = tag.add_run(
        "What the market is saying about GEO/AEO this week, "
        "drawn from practitioner forums, buyer reviews, and trade press."
    )
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    run.font.italic = True
    doc.add_paragraph("")

    # ── Helpers ──────────────────────────────────────────────────────────────
    recent_signals = sorted(
        [s for s in insights if s.get("post_date", "") >= _7d_ago],
        key=lambda s: _relevance_score(s), reverse=True,
    )
    if not recent_signals:
        recent_signals = sorted(
            insights, key=lambda s: _relevance_score(s), reverse=True,
        )[:20]

    def _sig_headline(s):
        return (s.get("title", "") or s.get("text", ""))[:120]

    def _sig_url(s):
        return s.get("url", "")

    # ── Section 1: Top Story ────────────────────────────────────────────────
    doc.add_heading("Top Story", level=1)

    # Use Signal of the Week if available, otherwise fall back to pipeline
    if _SOTW.get("title"):
        p = doc.add_paragraph()
        _docx_add_hyperlink(p, _SOTW["title"], _SOTW["url"], font_size=13, bold=True)

        doc.add_paragraph(_SOTW["brief"])

        p = doc.add_paragraph()
        r = p.add_run(f"{_SOTW['source']} | {_SOTW['date']}")
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    else:
        # Fallback: highest-relevance signal from the last 7 days
        top = recent_signals[0] if recent_signals else None
        if top:
            hl = _sig_headline(top)
            url = _sig_url(top)
            comps = top.get("companies_mentioned", [])
            sent = top.get("sentiment", "neutral")
            why = _relevance_sentence(top)

            p = doc.add_paragraph()
            if url:
                _docx_add_hyperlink(p, hl, url, font_size=13, bold=True)
            else:
                r = p.add_run(hl)
                r.bold = True
                r.font.size = Pt(13)

            comp_str = ", ".join(comps[:3]) if comps else "multiple players"
            source = top.get("source", "community forums")
            narrative = (
                f"This week's most significant signal comes from {source}, "
                f"involving {comp_str}. "
            )
            if why:
                narrative += why
            doc.add_paragraph(narrative)
        else:
            doc.add_paragraph("No high-confidence signals surfaced this week.")

    # ── Section 2: Who's Moving ─────────────────────────────────────────────
    doc.add_heading("Who's Moving", level=1)

    # Filter to companies with 3+ signals for the table
    qualified_comps = [
        c for c in selected_comps
        if comp_stats.get(c, {}).get("total", 0) >= 3
    ]
    qualified_comps.sort(
        key=lambda c: comp_stats.get(c, {}).get("total", 0), reverse=True,
    )
    table_comps = qualified_comps[:10]

    rising = [c for c in table_comps if comp_stats.get(c, {}).get("momentum") == "Rising"]
    falling = [c for c in table_comps if comp_stats.get(c, {}).get("momentum") == "Falling"]
    stable = [c for c in table_comps if comp_stats.get(c, {}).get("momentum") == "Stable"]

    # Narrative paragraph
    parts = []
    if rising:
        parts.append(
            f"{', '.join(rising[:4])} {'is' if len(rising) == 1 else 'are'} "
            f"gaining momentum with increased mention volume week over week"
        )
    if falling:
        parts.append(
            f"{', '.join(falling[:4])} {'is' if len(falling) == 1 else 'are'} "
            f"seeing declining conversation"
        )
    if stable:
        parts.append(
            f"{', '.join(stable[:4])} {'is' if len(stable) == 1 else 'are'} "
            f"holding steady"
        )
    if parts:
        doc.add_paragraph(". ".join(parts) + ".")
    else:
        doc.add_paragraph("Insufficient data to determine momentum trends this period.")

    # Momentum table (3+ signals only, max 10 rows)
    table = doc.add_table(rows=1, cols=5)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, h in enumerate(["Company", "Signals", "Positive %", "Momentum", "Last Signal"]):
        hdr[i].text = h
    for comp in table_comps:
        cs = comp_stats.get(comp, {})
        row = table.add_row().cells
        row[0].text = comp
        row[1].text = str(cs.get("total", 0))
        row[2].text = f"{cs.get('pos_pct', 0)}%"
        row[3].text = cs.get("momentum", "N/A")
        row[4].text = cs.get("latest", "N/A")
    _docx_source_caption(doc, total_signals, date_str)

    # ── Section 3: Voice of the Market ──────────────────────────────────────
    doc.add_heading("Voice of the Market", level=1)
    doc.add_paragraph(
        "The most notable signals from the past week, grouped by type. "
        "Buyer and feature-request signals are weighted highest."
    )

    # Buyer / feature request signals
    buyer_sigs = [
        s for s in recent_signals
        if s.get("is_buyer_voice") or s.get("is_feature_request")
    ]
    market_sigs = [
        s for s in recent_signals
        if s not in buyer_sigs and _is_displayable_post(s)
    ]

    def _render_signal_item(doc, s):
        """Render a signal as hyperlinked headline + metadata line."""
        hl = _sig_headline(s)
        url = _sig_url(s)
        src = _source_badge(s.get("source", ""))
        why = _relevance_sentence(s)
        comps = ", ".join(s.get("companies_mentioned", [])[:3])
        date = s.get("post_date", "")

        p = doc.add_paragraph(style="List Bullet")
        p.add_run(f"[{src}] ")
        if url:
            _docx_add_hyperlink(p, hl, url, font_size=11, bold=True)
        else:
            r = p.add_run(hl)
            r.bold = True

        # Metadata line: rationale | company | date
        meta_parts = []
        if why:
            meta_parts.append(why)
        if comps:
            meta_parts.append(comps)
        if date:
            meta_parts.append(date)
        if meta_parts:
            p2 = doc.add_paragraph()
            r2 = p2.add_run("  " + " | ".join(meta_parts))
            r2.font.size = Pt(9)
            r2.font.italic = True
            r2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    if buyer_sigs:
        doc.add_heading("Buyer Signals", level=2)
        for s in buyer_sigs[:5]:
            _render_signal_item(doc, s)

    if market_sigs:
        doc.add_heading("Market Signals", level=2)
        for s in market_sigs[:5]:
            _render_signal_item(doc, s)

    if not buyer_sigs and not market_sigs:
        doc.add_paragraph("No notable signals this period.")

    # ── Section 4: What Buyers Are Asking For ───────────────────────────────
    doc.add_heading("What Buyers Are Asking For", level=1)
    doc.add_paragraph(
        "The feature themes generating the most conversation, "
        "ranked by evidence volume and recency."
    )

    sorted_feats = sorted(
        opportunity_data.items(),
        key=lambda x: x[1]["evidence"],
        reverse=True,
    )
    top_feats = [(f, od) for f, od in sorted_feats if od["evidence"] >= 2][:3]

    for feat, od in top_feats:
        doc.add_heading(feat, level=2)

        # Analyst-style narrative
        recent_ct = sum(1 for s in od["signals"] if s.get("post_date", "") >= _90d_ago)
        praised = list(od["companies_praised"])[:3]
        complained = list(od["companies_complained"])[:3]
        requests_ct = od["requests"]

        narrative_parts = [
            f"{feat} has {od['evidence']} total evidence signals "
            f"({od['confidence']}% confidence*)."
        ]
        if recent_ct:
            narrative_parts.append(
                f"{recent_ct} of those appeared in the last 90 days."
            )
        if requests_ct:
            narrative_parts.append(
                f"{requests_ct} are explicit feature requests from users."
            )
        if praised:
            narrative_parts.append(
                f"Companies receiving praise here: {', '.join(praised)}."
            )
        if complained:
            narrative_parts.append(
                f"Companies drawing complaints: {', '.join(complained)}."
            )
        doc.add_paragraph(" ".join(narrative_parts))

        # Top evidence signals with hyperlinked headlines
        top_ev = sorted(od["signals"], key=lambda s: _relevance_score(s), reverse=True)[:3]
        for s in top_ev:
            _render_signal_item(doc, s)

    if not top_feats:
        doc.add_paragraph("No features have enough evidence to highlight this period.")
    else:
        _docx_confidence_footnote(doc)

    # ── Section 5: Ones to Watch ────────────────────────────────────────────
    doc.add_heading("Ones to Watch", level=1)
    doc.add_paragraph(
        "Companies that are new to the dataset or showing unusual "
        "activity patterns worth monitoring."
    )

    # Find rising or newly-appeared companies
    watch_candidates = []
    median_total = sorted(
        [comp_stats.get(c, {}).get("total", 0) for c in selected_comps]
    )[len(selected_comps) // 2] if selected_comps else 5

    for comp in selected_comps:
        cs = comp_stats.get(comp, {})
        total = cs.get("total", 0)
        mom = cs.get("momentum", "")
        if mom == "Rising" and total >= 2:
            watch_candidates.append((comp, cs))

    try:
        new_cos = _get_new_companies(json.dumps(insights))
    except Exception:
        new_cos = set()
    for comp in selected_comps:
        if comp in new_cos and comp not in [w[0] for w in watch_candidates]:
            cs = comp_stats.get(comp, {})
            if cs.get("total", 0) >= 2:
                watch_candidates.append((comp, cs))

    # Sort: below-median rising first (emerging), then above-median rising
    watch_candidates.sort(
        key=lambda x: (0 if x[1].get("total", 0) <= median_total else 1, -x[1].get("total", 0))
    )
    watch_candidates = watch_candidates[:3]

    if watch_candidates:
        for comp, cs in watch_candidates:
            meta = company_meta.get(comp, {})
            positioning = _company_positioning(meta) if meta else ""
            total = cs.get("total", 0)
            pos_pct = cs.get("pos_pct", 0)
            mom = cs.get("momentum", "N/A")
            is_new = comp in new_cos

            p = doc.add_paragraph()
            r = p.add_run(comp)
            r.bold = True
            r.font.size = Pt(12)

            # Build analyst narrative
            comp_sigs = [
                s for s in insights
                if comp in s.get("companies_mentioned", [])
                and _is_displayable_post(s)
            ]
            comp_sigs.sort(key=lambda s: _relevance_score(s), reverse=True)
            best = comp_sigs[0] if comp_sigs else None
            best_why = _relevance_sentence(best, for_company=comp) if best else ""

            sentences = []
            if positioning and total > median_total:
                sentences.append(
                    f"{comp} ({positioning}) is one of the most active competitors "
                    f"in the dataset this week with {total} signals."
                )
            elif positioning and is_new:
                sentences.append(
                    f"{comp} ({positioning}) is a new entrant, first appearing in "
                    f"the data within the last seven days with {total} signals."
                )
            elif positioning:
                sentences.append(
                    f"{comp} ({positioning}) is gaining momentum with {total} "
                    f"signals and {pos_pct}% positive sentiment."
                )
            else:
                sentences.append(
                    f"{comp} is showing rising activity with {total} signals "
                    f"tracked so far."
                )

            if best_why:
                sentences.append(f"The top signal: {best_why}")

            if best:
                best_src = best.get("source", "")
                best_hl = _sig_headline(best)
                best_url = _sig_url(best)
                if best_url and best_src:
                    sentences.append(
                        f"Most notable mention came via {best_src}."
                    )

            doc.add_paragraph(" ".join(sentences))

            # Hyperlinked top signal
            if best:
                best_hl = _sig_headline(best)
                best_url = _sig_url(best)
                best_src = _source_badge(best.get("source", ""))
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(f"[{best_src}] ")
                if best_url:
                    _docx_add_hyperlink(p, best_hl, best_url, font_size=11, bold=False)
                else:
                    p.add_run(best_hl)
    else:
        doc.add_paragraph("No emerging players flagged this period.")

    # ── Methodology footnote ────────────────────────────────────────────────
    doc.add_paragraph("")
    p = doc.add_paragraph()
    run = p.add_run(_METHODOLOGY_NOTE)
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    p = doc.add_paragraph()
    run = p.add_run(
        f"Sources: {SOURCE_LIST}. "
        f"Pipeline refreshes every 6 hours. {total_signals:,} signals in database. "
        f"Report generated {date_display}."
    )
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _export_briefing_deck(insights, company_meta, opportunity_data, selected_comps, comp_stats, fig1_bytes, fig2_bytes):
    """Generate a Briefing Deck .pptx and return bytes."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    total_signals = len(insights)
    date_str = datetime.now().strftime("%Y-%m-%d")
    date_display = datetime.now().strftime("%B %d, %Y")

    # Pre-compute callout data
    rising = [c for c in selected_comps if comp_stats.get(c, {}).get("momentum") == "Rising"]
    falling = [c for c in selected_comps if comp_stats.get(c, {}).get("momentum") == "Falling"]
    mean_signals = sum(
        comp_stats.get(c, {}).get("total", 0) for c in selected_comps
    ) / max(len(selected_comps), 1)
    leader = max(selected_comps, key=lambda c: comp_stats.get(c, {}).get("total", 0)) if selected_comps else ""
    leader_total = comp_stats.get(leader, {}).get("total", 0) if leader else 0
    above_avg = [c for c in selected_comps if comp_stats.get(c, {}).get("total", 0) > mean_signals]
    below_avg = [c for c in selected_comps if comp_stats.get(c, {}).get("total", 0) < mean_signals]

    # Heat map callout data
    feat_totals = {}
    best_cell_score = 0
    best_cell_comp = ""
    best_cell_feat = ""
    for opp, od in opportunity_data.items():
        feat_sum = sum(
            od["company_detail"].get(c, {}).get("count", 0) if hasattr(od["company_detail"], "get") else 0
            for c in selected_comps
        )
        feat_totals[opp] = feat_sum
        for c in selected_comps:
            cd = od["company_detail"].get(c, {}) if hasattr(od["company_detail"], "get") else {}
            ct = cd.get("count", 0) if isinstance(cd, dict) else 0
            if ct > best_cell_score:
                best_cell_score = ct
                best_cell_comp = c
                best_cell_feat = opp
    hottest_feat = max(feat_totals, key=feat_totals.get) if feat_totals else ""

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.5), PptxInches(11), PptxInches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "GEO Pulse Briefing"
    p.font.size = PptxPt(44)
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = f"Market Intelligence Report | {date_display}"
    p2.font.size = PptxPt(20)
    p2.font.color.rgb = PptxRGBColor(0x64, 0x64, 0x64)

    # Slide 2: Snapshot
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(1))
    tf = txBox.text_frame
    tf.paragraphs[0].text = "Market Snapshot"
    tf.paragraphs[0].font.size = PptxPt(32)
    tf.paragraphs[0].font.bold = True

    stats_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(5))
    tf2 = stats_box.text_frame
    tf2.word_wrap = True
    lines = [
        f"Total Signals: {total_signals:,}",
        f"Companies Tracked: {len(company_meta)}",
        f"Competitors in View: {len(selected_comps)}",
        "",
    ]
    if rising:
        lines.append(f"Rising: {', '.join(rising)}")
    if falling:
        lines.append(f"Falling: {', '.join(falling)}")
    tf2.paragraphs[0].text = lines[0]
    tf2.paragraphs[0].font.size = PptxPt(18)
    for line in lines[1:]:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = PptxPt(18)
    _pptx_slide_footer(slide, date_str, total_signals)

    # Slide 3: Momentum Chart
    if fig1_bytes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.6))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "Competitor Presence & Momentum"
        tf.paragraphs[0].font.size = PptxPt(28)
        tf.paragraphs[0].font.bold = True
        img_stream = io.BytesIO(fig1_bytes)
        slide.shapes.add_picture(img_stream, PptxInches(0.5), PptxInches(1.0), PptxInches(12), PptxInches(4.2))

        # Auto-generated callouts
        momentum_bullets = []
        if leader:
            momentum_bullets.append(
                f"\u2022 {leader} leads in signal volume with {leader_total} mentions"
            )
        if rising:
            momentum_bullets.append(
                f"\u2022 {', '.join(rising[:3])} {'is' if len(rising) == 1 else 'are'} "
                f"gaining momentum week over week"
            )
        above_str = ", ".join(above_avg[:3]) if above_avg else "none"
        momentum_bullets.append(
            f"\u2022 Category average is {mean_signals:.0f} signals. "
            f"{above_str} stand{'s' if len(above_avg) == 1 else ''} out above average"
        )
        _pptx_callout_box(slide, momentum_bullets, top_inches=5.3)
        _pptx_slide_footer(slide, date_str, total_signals)

    # Slide 4: Heat Map
    if fig2_bytes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.6))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "Feature Heat Map"
        tf.paragraphs[0].font.size = PptxPt(28)
        tf.paragraphs[0].font.bold = True
        img_stream = io.BytesIO(fig2_bytes)
        slide.shapes.add_picture(img_stream, PptxInches(0.5), PptxInches(1.0), PptxInches(12), PptxInches(4.2))

        # Auto-generated callouts
        heat_bullets = []
        if hottest_feat:
            heat_bullets.append(
                f"\u2022 {hottest_feat} shows the most market activity across competitors"
            )
        if best_cell_comp and best_cell_feat:
            heat_bullets.append(
                f"\u2022 {best_cell_comp} has the strongest recent coverage in {best_cell_feat}"
            )
        heat_bullets.append(
            "\u2022 White cells = no public signals found, not confirmed absence of capability"
        )
        _pptx_callout_box(slide, heat_bullets, top_inches=5.3)
        _pptx_slide_footer(slide, date_str, total_signals)

    # Slide 5: Build Now
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(1))
    tf = txBox.text_frame
    tf.paragraphs[0].text = "Build Now Opportunities"
    tf.paragraphs[0].font.size = PptxPt(28)
    tf.paragraphs[0].font.bold = True
    bn_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(4.5))
    tf2 = bn_box.text_frame
    tf2.word_wrap = True
    bn_count = 0
    for opp, od in sorted(opportunity_data.items(), key=lambda x: x[1]["evidence"], reverse=True):
        if od["evidence"] < 3:
            continue
        red_count = sum(
            1 for c in selected_comps
            if c not in od["companies_praised"]
            and c not in od["companies_complained"]
            and c not in od["companies_tried"]
        )
        if red_count > len(selected_comps) / 2:
            p = tf2.add_paragraph() if bn_count > 0 else tf2.paragraphs[0]
            p.text = f"{opp} ({od['confidence']}% confidence*, {od['evidence']} signals)"
            p.font.size = PptxPt(16)
            bn_count += 1
    if bn_count == 0:
        tf2.paragraphs[0].text = "No features currently meet Build Now criteria."
        tf2.paragraphs[0].font.size = PptxPt(16)
    # Confidence footnote
    fn_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(6.3), PptxInches(12), PptxInches(0.5))
    fn_tf = fn_box.text_frame
    fn_tf.paragraphs[0].text = f"* {_CONFIDENCE_FOOTNOTE}"
    fn_tf.paragraphs[0].font.size = PptxPt(9)
    fn_tf.paragraphs[0].font.color.rgb = PptxRGBColor(0x99, 0x99, 0x99)
    fn_tf.paragraphs[0].font.italic = True
    _pptx_slide_footer(slide, date_str, total_signals)

    # Slide 6: About This Data
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(1))
    tf = txBox.text_frame
    tf.paragraphs[0].text = "About This Data"
    tf.paragraphs[0].font.size = PptxPt(28)
    tf.paragraphs[0].font.bold = True

    meth_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.3), PptxInches(12), PptxInches(2.5))
    meth_tf = meth_box.text_frame
    meth_tf.word_wrap = True
    meth_tf.paragraphs[0].text = _METHODOLOGY_NOTE
    meth_tf.paragraphs[0].font.size = PptxPt(13)

    src_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(3.8), PptxInches(12), PptxInches(3.2))
    src_tf = src_box.text_frame
    src_tf.word_wrap = True
    p0 = src_tf.paragraphs[0]
    p0.text = "Sources Monitored"
    p0.font.size = PptxPt(16)
    p0.font.bold = True
    for name, desc in _SOURCE_DESCRIPTIONS:
        p = src_tf.add_paragraph()
        p.text = f"\u2022 {name}, {desc}"
        p.font.size = PptxPt(11)
        p.font.color.rgb = PptxRGBColor(0x33, 0x33, 0x33)

    p_total = src_tf.add_paragraph()
    p_total.text = f"\nTotal signals in database: {total_signals:,}. Pipeline refreshes every 6 hours."
    p_total.font.size = PptxPt(10)
    p_total.font.color.rgb = PptxRGBColor(0x66, 0x66, 0x66)
    _pptx_slide_footer(slide, date_str, total_signals)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _export_prd(opportunity_data, insights, selected_features, selected_comps):
    """Generate a PRD .docx for selected features and return bytes."""
    doc = DocxDocument()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    total_signals = len(insights)
    date_str = datetime.now().strftime("%Y-%m-%d")
    date_display = datetime.now().strftime("%B %d, %Y")

    doc.add_heading("Product Requirements Document", level=0)
    doc.add_paragraph(f"Generated {date_display} by GEO Pulse")
    # Cover page footer
    p = doc.add_paragraph()
    run = p.add_run(
        "Market data sourced from GEO Pulse. All signals are public. "
        "Original sources linked in Signal Appendix."
    )
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.italic = True
    doc.add_paragraph("")

    _90d = (datetime.now() - timedelta(days=90)).strftime("%Y-%m-%d")

    for feat in selected_features:
        od = opportunity_data.get(feat)
        if not od:
            continue

        doc.add_heading(feat, level=1)

        # Overview
        doc.add_heading("Overview", level=2)
        doc.add_paragraph(
            f"This feature has {od['evidence']} evidence signals with "
            f"{od['confidence']}% confidence score.*"
        )
        _docx_confidence_footnote(doc)

        # Problem Statement
        doc.add_heading("Problem Statement", level=2)
        complaints = [s for s in od["signals"] if "complaint" in s.get("entity_tags", [])]
        if complaints:
            doc.add_paragraph(
                f"{len(complaints)} complaint signals identified. Key themes:"
            )
            for s in complaints[:5]:
                stitle = s.get("title", "")[:100] or s.get("text", "")[:100]
                doc.add_paragraph(f"- {stitle}", style="List Bullet")
        else:
            doc.add_paragraph("No direct complaints found. Demand driven by feature requests and market gaps.")

        # Proposed Solution
        doc.add_heading("Proposed Solution", level=2)
        doc.add_paragraph("[To be completed by product team]")

        # Success Metrics
        doc.add_heading("Success Metrics", level=2)
        doc.add_paragraph("[To be completed by product team]")

        # Competitive Landscape
        doc.add_heading("Competitive Landscape", level=2)
        praised = od["companies_praised"] & set(selected_comps) if selected_comps else od["companies_praised"]
        complained = od["companies_complained"] & set(selected_comps) if selected_comps else od["companies_complained"]
        no_data = [c for c in selected_comps if c not in od["companies_praised"]
                   and c not in od["companies_complained"] and c not in od["companies_tried"]]

        if praised:
            doc.add_paragraph(f"Praised: {', '.join(praised)}")
        if complained:
            doc.add_paragraph(f"Complaints: {', '.join(complained)}")
        if no_data:
            doc.add_paragraph(f"No data: {', '.join(no_data)}")

        # Open Questions
        doc.add_heading("Open Questions", level=2)
        doc.add_paragraph("[To be completed by product team]")

        # Signal Appendix (expanded with full citations)
        doc.add_heading("Signal Appendix", level=2)
        recent = [s for s in od["signals"] if s.get("post_date", "") >= _90d]
        display_sigs = recent[:20] if recent else od["signals"][:20]
        table = doc.add_table(rows=1, cols=6)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        for i, h in enumerate(["Title", "Source", "URL", "Date", "Sentiment", "Rationale"]):
            hdr[i].text = h
        for s in display_sigs:
            row = table.add_row().cells
            row[0].text = (s.get("title", "") or s.get("text", ""))[:80]
            row[1].text = s.get("source", "")
            row[2].text = s.get("url", "")
            row[3].text = s.get("post_date", "")
            row[4].text = s.get("sentiment", "")
            row[5].text = _relevance_sentence(s) or ""
        _docx_source_caption(doc, total_signals, date_str)

        doc.add_page_break()

    # Document-level methodology
    _docx_methodology_appendix(doc, total_signals)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _export_brd(opportunity_data, insights, selected_features, selected_comps):
    """Generate a BRD .docx for selected features and return bytes."""
    doc = DocxDocument()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    total_signals = len(insights)
    date_str = datetime.now().strftime("%Y-%m-%d")
    date_display = datetime.now().strftime("%B %d, %Y")

    doc.add_heading("Business Requirements Document", level=0)
    doc.add_paragraph(f"Generated {date_display} by GEO Pulse")
    # Cover page footer
    p = doc.add_paragraph()
    run = p.add_run(
        "Market data sourced from GEO Pulse. All signals are public. "
        "Original sources linked in Signal Appendix."
    )
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.italic = True
    doc.add_paragraph("")

    _90d = (datetime.now() - timedelta(days=90)).strftime("%Y-%m-%d")

    for feat in selected_features:
        od = opportunity_data.get(feat)
        if not od:
            continue

        doc.add_heading(feat, level=1)

        # Executive Summary
        doc.add_heading("Executive Summary", level=2)
        doc.add_paragraph(
            f"Market evidence supports investment in {feat}. "
            f"{od['evidence']} signals identified with {od['confidence']}% confidence.*"
        )
        _docx_confidence_footnote(doc)

        # Business Objective
        doc.add_heading("Business Objective", level=2)
        doc.add_paragraph("[To be completed by business stakeholders]")

        # Market Evidence
        doc.add_heading("Market Evidence", level=2)
        doc.add_paragraph(f"Total signals: {od['evidence']}")
        doc.add_paragraph(f"Complaints: {od['complaints']}")
        doc.add_paragraph(f"Feature requests: {od['requests']}")
        doc.add_paragraph(f"Praise: {od['praise']}")
        recent_ct = sum(1 for s in od["signals"] if s.get("post_date", "") >= _90d)
        doc.add_paragraph(f"Signals in last 90 days: {recent_ct}")

        # Stakeholders
        doc.add_heading("Stakeholders", level=2)
        doc.add_paragraph("[To be completed by business stakeholders]")

        # Scope
        doc.add_heading("Scope", level=2)
        doc.add_paragraph("[To be completed by business stakeholders]")

        # Constraints and Assumptions
        doc.add_heading("Constraints and Assumptions", level=2)
        doc.add_paragraph("[To be completed by business stakeholders]")

        # Approval
        doc.add_heading("Approval", level=2)
        doc.add_paragraph("[Pending approval]")

        # Signal Appendix (expanded with full citations)
        doc.add_heading("Signal Appendix", level=2)
        recent = [s for s in od["signals"] if s.get("post_date", "") >= _90d]
        display_sigs = recent[:20] if recent else od["signals"][:20]
        table = doc.add_table(rows=1, cols=6)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        for i, h in enumerate(["Title", "Source", "URL", "Date", "Sentiment", "Rationale"]):
            hdr[i].text = h
        for s in display_sigs:
            row = table.add_row().cells
            row[0].text = (s.get("title", "") or s.get("text", ""))[:80]
            row[1].text = s.get("source", "")
            row[2].text = s.get("url", "")
            row[3].text = s.get("post_date", "")
            row[4].text = s.get("sentiment", "")
            row[5].text = _relevance_sentence(s) or ""
        _docx_source_caption(doc, total_signals, date_str)

        doc.add_page_break()

    # Document-level methodology
    _docx_methodology_appendix(doc, total_signals)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Email digest (SendGrid)
# ---------------------------------------------------------------------------

SUBSCRIBERS_PATH = os.path.join(DATA_DIR, "subscribers.json")
EMAIL_LOG_PATH = os.path.join(DATA_DIR, "email_log.json")


def _load_subscribers():
    if os.path.exists(SUBSCRIBERS_PATH):
        with open(SUBSCRIBERS_PATH, "r") as f:
            return json.load(f)
    return []


def _save_subscribers(subs):
    with open(SUBSCRIBERS_PATH, "w") as f:
        json.dump(subs, f, ensure_ascii=False, indent=2)


def _log_email(entry):
    log = []
    if os.path.exists(EMAIL_LOG_PATH):
        with open(EMAIL_LOG_PATH, "r") as f:
            log = json.load(f)
    log.append(entry)
    # Keep last 500 entries
    with open(EMAIL_LOG_PATH, "w") as f:
        json.dump(log[-500:], f, ensure_ascii=False, indent=2)


def _send_email(to_email, subject, html_content):
    """Send an email via SendGrid. Returns True on success."""
    try:
        import sendgrid
        from sendgrid.helpers.mail import Mail, Email, To, Content
    except ImportError:
        return False

    api_key = os.environ.get("SENDGRID_API_KEY") or st.secrets.get("SENDGRID_API_KEY", "")
    from_email = os.environ.get("SENDGRID_FROM_EMAIL") or st.secrets.get("SENDGRID_FROM_EMAIL", "noreply@geopulse.io")

    if not api_key:
        return False

    sg = sendgrid.SendGridAPIClient(api_key=api_key)
    message = Mail(
        from_email=Email(from_email),
        to_emails=To(to_email),
        subject=subject,
        html_content=Content("text/html", html_content),
    )
    try:
        response = sg.send(message)
        return response.status_code in (200, 201, 202)
    except Exception:
        return False


def _send_confirmation_email(email, name):
    """Send a subscription confirmation email."""
    html = f"""
    <div style="font-family: 'DM Sans', Arial, sans-serif; max-width: 600px; margin: 0 auto;
    background: #F8F4EB; color: #0A0A0A; padding: 24px;">
        <h2 style="color: #0A0A0A;">Welcome to GEO Pulse Daily Digest</h2>
        <p>Hi {name or 'there'},</p>
        <p>You're now subscribed to the GEO Pulse daily digest.
        You'll receive market intelligence updates at your preferred time.</p>
        <p style="color: #0A0A0A; font-size: 0.85rem;">
            To unsubscribe, visit the GEO Pulse dashboard and remove your subscription
            in the Email Digest section.</p>
        <hr style="border: none; border-top: 1px solid #D1CFBA;">
        <p style="color: #D1CFBA; font-size: 0.75rem;">GEO Pulse Market Intelligence</p>
    </div>
    """
    return _send_email(email, "Welcome to GEO Pulse Daily Digest", html)


def _build_digest_html(insights_list, company_meta_dict, comp_filter=None):
    """Build the daily digest HTML email content."""
    now = datetime.now()
    week_ago = (now - timedelta(days=7)).strftime("%Y-%m-%d")
    yesterday = (now - timedelta(days=1)).strftime("%Y-%m-%d")

    # Filter by company if specified
    if comp_filter:
        relevant = [i for i in insights_list
                    if any(c in i.get("companies_mentioned", []) for c in comp_filter)]
    else:
        relevant = insights_list

    recent = [i for i in relevant if i.get("post_date", "") >= yesterday]
    if not recent:
        recent = [i for i in relevant if i.get("post_date", "") >= week_ago]

    # Sort by relevance
    recent.sort(key=lambda x: _relevance_score(x), reverse=True)
    top5 = recent[:5]

    # Momentum snapshot
    comp_counts = Counter()
    for i in relevant:
        for c in i.get("companies_mentioned", []):
            comp_counts[c] += 1
    top_comps = comp_counts.most_common(5)

    # Build HTML
    signal_rows = ""
    for s in top5:
        title = s.get("title", "")[:100] or s.get("text", "")[:100]
        url = s.get("url", "")
        source = s.get("source", "")
        date = s.get("post_date", "")
        link = f'<a href="{url}" style="color: #0E3B7E; text-decoration: none;">{title}</a>' if url else title
        signal_rows += f"""
        <tr>
            <td style="padding: 8px; border-bottom: 1px solid #D1CFBA;">
                <span style="background: #D1CFBA; padding: 2px 6px;
                font-family: monospace; font-size: 0.75rem;">{source}</span>
                {link}
                <br><span style="color: #888; font-size: 0.8rem;">{date}</span>
            </td>
        </tr>
        """

    momentum_rows = ""
    for comp, count in top_comps:
        momentum_rows += f"""
        <tr>
            <td style="padding: 4px 8px; border-bottom: 1px solid #D1CFBA;">{comp}</td>
            <td style="padding: 4px 8px; border-bottom: 1px solid #D1CFBA;">{count} signals</td>
        </tr>
        """

    html = f"""
    <div style="font-family: 'DM Sans', Arial, sans-serif; max-width: 650px; margin: 0 auto;
    background: #F8F4EB; border: 1px solid #D1CFBA;">
        <div style="background: #0A0A0A; color: #F8F4EB; padding: 16px 24px;">
            <h1 style="margin: 0; font-size: 1.3rem;">GEO Pulse Daily Digest</h1>
            <p style="margin: 4px 0 0 0; color: #D1CFBA; font-size: 0.85rem;">
                {now.strftime('%B %d, %Y')}</p>
        </div>

        <div style="padding: 20px 24px;">
            <h2 style="font-size: 1.1rem; color: #0A0A0A; margin-top: 0;">Momentum Snapshot</h2>
            <table style="width: 100%; border-collapse: collapse;">
                {momentum_rows}
            </table>

            <h2 style="font-size: 1.1rem; color: #0A0A0A; margin-top: 24px;">Top Signals</h2>
            <table style="width: 100%; border-collapse: collapse;">
                {signal_rows}
            </table>

            <div style="margin-top: 24px; padding: 12px; background: #D1CFBA;
            font-size: 0.85rem; color: #0A0A0A;">
                Open the <a href="#" style="color: #0E3B7E;">GEO Pulse dashboard</a>
                for full details, charts, and export options.
            </div>
        </div>

        <div style="padding: 12px 24px; background: #D1CFBA;
        font-size: 0.75rem; color: #0A0A0A;">
            GEO Pulse Market Intelligence | To unsubscribe, visit the dashboard settings.
        </div>
    </div>
    """
    return html


def _send_daily_digests():
    """Scheduled job: send daily digest emails to all subscribers."""
    subs = _load_subscribers()
    if not subs:
        return

    # Load fresh data
    if os.path.exists(INSIGHTS_PATH):
        with open(INSIGHTS_PATH, "r") as f:
            all_insights = json.load(f)
    else:
        return

    now = datetime.now()
    current_hour = now.hour

    for sub in subs:
        if not sub.get("confirmed", False):
            continue

        # Check delivery time preference
        pref_hour = sub.get("delivery_hour", 8)
        tz_offset = sub.get("tz_offset", 0)
        adjusted_hour = (current_hour + tz_offset) % 24

        if adjusted_hour != pref_hour:
            continue

        comp_filter = sub.get("competitor_filter", [])
        html = _build_digest_html(all_insights, {}, comp_filter or None)

        ok = _send_email(sub["email"], "GEO Pulse Daily Digest", html)
        _log_email({
            "email": sub["email"],
            "sent_at": now.isoformat(),
            "success": ok,
            "type": "daily_digest",
        })


# ---------------------------------------------------------------------------
# Display-level relevance and age filters
# ---------------------------------------------------------------------------

_GEO_DISPLAY_TERMS = [
    "geo ", "geo/aeo", " aeo ", "generative engine", "answer engine",
    "ai search", "ai visibility", "ai answer", "ai citation", "ai overview",
    "brand visibility", "share of voice", "share of answer",
    "llm optimization", "llm brand", "llm monitoring",
    "ai overviews", "zero click", "zero-click",
    "content optimization", "structured data", "schema markup",
    "searchgpt", "search gpt",
    "chatgpt search", "chatgpt visibility", "chatgpt citation",
    "gemini search", "gemini visibility",
    "perplexity search", "perplexity answer", "perplexity citation",
    "seo tool", "seo platform", "seo measurement",
    "brand mention", "brand monitoring", "brand measurement",
    "generative search", "conversational search",
]

# Weak terms require a company mention to qualify
_GEO_WEAK_TERMS = [
    "seo", "chatgpt", "perplexity", "gemini", "ai tool",
]

_MAX_AGE_DAYS = 730  # 24 months


def _is_display_relevant(insight):
    """Require GEO context for display. Company mention alone is insufficient."""
    text = (insight.get("text", "") + " " + insight.get("title", "")).lower()
    title = (insight.get("title") or "").strip()

    # Block clearly off-topic content
    if _OFFTOPIC_BLOCKLIST.search(text):
        return False
    if _TITLE_BLOCKLIST.search(title):
        return False

    # Require GEO/AEO context keywords
    has_context = any(term in text for term in _GEO_DISPLAY_TERMS)
    if has_context:
        return True
    # Weak terms only pass with a company mention
    has_companies = bool(insight.get("companies_mentioned"))
    has_weak = any(term in text for term in _GEO_WEAK_TERMS)
    if has_companies and has_weak:
        return True
    source = insight.get("source", "")
    if has_companies and source in ("G2", "Product Hunt"):
        return True
    return False


def _within_age_limit(insight):
    """Exclude articles older than 24 months."""
    try:
        dt = datetime.strptime(insight.get("post_date", ""), "%Y-%m-%d")
        cutoff = datetime.now() - timedelta(days=_MAX_AGE_DAYS)
        return dt >= cutoff
    except (ValueError, TypeError):
        return False


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Startup staleness check — refresh pipeline if data is old
# ---------------------------------------------------------------------------

def _data_is_stale(threshold_hours=6):
    """Check if last pipeline run was more than threshold_hours ago.
    Uses run_log.json on disk (survives hibernation) and falls back to
    file mtime of enriched_insights.json."""
    try:
        if os.path.exists(RUN_LOG_PATH):
            with open(RUN_LOG_PATH, "r") as f:
                log = json.load(f)
            if log:
                last_ts = log[-1].get("completed_at", "")
                if last_ts:
                    last_dt = datetime.fromisoformat(last_ts)
                    return (datetime.now() - last_dt).total_seconds() > threshold_hours * 3600
        # Fallback: check file modification time
        if os.path.exists(INSIGHTS_PATH):
            mtime = os.path.getmtime(INSIGHTS_PATH)
            age_hours = (datetime.now().timestamp() - mtime) / 3600
            return age_hours > threshold_hours
    except Exception:
        pass
    return True  # No data at all, definitely stale


if _data_is_stale():
    try:
        run_enrichment()
        st.cache_data.clear()
    except Exception:
        pass  # Pipeline failure is non-fatal; app still loads existing data

_raw_insights = load_insights()
insights = _dedup_insights([i for i in _raw_insights
                            if _within_age_limit(i) and _is_displayable_post(i)])
trends = load_trends()
companies_data = load_companies()
discovered = load_discovered_sources()

own_brands = {c["name"] for c in companies_data.get("own_brands", [])}
company_meta = {}
for group in ("own_brands", "competitors"):
    for c in companies_data.get(group, []):
        company_meta[c["name"]] = c

# Source stats for header
skip_domains = {"preview.redd.it", "i.redd.it", "v.redd.it", "sh.reddit.com", "imgur.com"}
approved_sources = [s for s in discovered
                    if s.get("status") == "approved" and s.get("domain", "") not in skip_domains]

# ---------------------------------------------------------------------------
# Sidebar: Ask GEO Pulse
# ---------------------------------------------------------------------------

def _build_data_summary():
    """Build a compact data summary for the LLM context."""
    from collections import Counter as _C
    comp_counts = _C()
    comp_sentiment = defaultdict(lambda: {"pos": 0, "neg": 0, "total": 0})
    tag_counts = _C()
    source_counts = _C()
    week_ago = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d")
    this_week = []

    for i in insights:
        for c in i.get("companies_mentioned", []):
            comp_counts[c] += 1
            s = i.get("sentiment", "neutral")
            comp_sentiment[c]["total"] += 1
            if s == "positive":
                comp_sentiment[c]["pos"] += 1
            elif s == "negative":
                comp_sentiment[c]["neg"] += 1
        for t in i.get("entity_tags", []):
            tag_counts[t] += 1
        source_counts[i.get("source", "")] += 1
        if i.get("post_date", "") >= week_ago:
            this_week.append(i)

    # Top complaints this week
    complaints = [i for i in this_week if "complaint" in i.get("entity_tags", [])]
    requests = [i for i in this_week if i.get("is_feature_request")]

    # Rising companies from trends
    rising = []
    fading = []
    if trends:
        for r in trends.get("rising", [])[:5]:
            rising.append(f"{r['name']} (+{r['delta']}%)")
        for r in trends.get("fading", [])[:3]:
            fading.append(f"{r['name']} ({r['delta']}%)")

    summary = f"""DATA SUMMARY (as of {datetime.now().strftime('%Y-%m-%d %H:%M')}):
- Total signals: {len(insights):,} from {len(source_counts)} source types
- This week: {len(this_week)} signals, {len(complaints)} complaints, {len(requests)} feature requests
- Companies tracked: {len(comp_counts)}
- Top mentioned: {', '.join(f'{c} ({n})' for c, n in comp_counts.most_common(10))}
- Sources: {', '.join(f'{s}: {n}' for s, n in source_counts.most_common())}
- Rising: {', '.join(rising) if rising else 'None'}
- Fading: {', '.join(fading) if fading else 'None'}
- Top complaints: {dict(tag_counts)}

SENTIMENT BY COMPANY (top 15):
"""
    for c, _ in comp_counts.most_common(15):
        cs = comp_sentiment[c]
        pos_pct = round(cs["pos"] * 100 / max(cs["total"], 1))
        neg_pct = round(cs["neg"] * 100 / max(cs["total"], 1))
        summary += f"- {c}: {cs['total']} mentions, {pos_pct}% positive, {neg_pct}% negative\n"

    return summary


def _get_starter_questions(active_tab_idx=0):
    """Generate 5 starter questions based on current data and active tab."""
    # Base questions that always work
    base = [
        "What are buyers complaining about most this week?",
        "Which company is gaining the most momentum right now?",
        "What features should we prioritize building next?",
    ]

    # Data-driven additions
    if trends:
        rising = trends.get("rising", [])
        for r in rising:
            if r.get("name") in company_meta:
                base.append(f"Why is {r['name']} gaining momentum?")
                break

    comp_counts = Counter()
    for i in insights:
        for c in i.get("companies_mentioned", []):
            comp_counts[c] += 1
    top2 = [c for c, _ in comp_counts.most_common(2)]
    if len(top2) >= 2:
        base.append(f"What are people saying about {top2[0]} vs {top2[1]}?")

    if not any("gap" in q.lower() for q in base):
        base.append("What's the biggest gap no tool has solved yet?")

    # Tab-specific biasing
    if active_tab_idx == 1:  # Sentiment
        base.insert(0, "Which tools have the most negative sentiment and why?")
    elif active_tab_idx == 4:  # Roadmap
        base.insert(0, "What product opportunities have the strongest evidence?")
    elif active_tab_idx == 0:  # Live Feed
        base.insert(0, "What's the most important signal from the last 48 hours?")

    return base[:5]


def _get_relevant_posts(query, limit=15):
    """Find the most relevant posts for a given query."""
    query_lower = query.lower()
    scored = []
    for i in insights:
        text = (i.get("text", "") + " " + i.get("title", "")).lower()
        score = 0
        for word in query_lower.split():
            if len(word) > 3 and word in text:
                score += 1
        # Boost for company name matches
        for c in i.get("companies_mentioned", []):
            if c.lower() in query_lower:
                score += 3
        if score > 0:
            scored.append((score, i))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [s[1] for s in scored[:limit]]


_ASK_AI_RENDERED_ABOVE_TABS = True  # Ask GEO Pulse lives above the tab bar


# ---------------------------------------------------------------------------
# Header
# ---------------------------------------------------------------------------

st.title("GEO Pulse")
st.caption("Market intelligence for the GEO/AEO category, so your team always knows before the competition.")

# Persistent header bar
runs = load_run_log()
last_run = runs[-1] if runs else {}
last_ts = last_run.get("completed_at", "")
try:
    last_dt = datetime.fromisoformat(last_ts)
    _delta_sec = (datetime.now() - last_dt).total_seconds()
    if _delta_sec < 60:
        freshness = "just now"
    elif _delta_sec < 3600:
        freshness = f"{int(_delta_sec // 60)}m ago"
    elif _delta_sec < 86400:
        freshness = f"{int(_delta_sec // 3600)}h ago"
    else:
        freshness = f"{_delta_sec / 86400:.1f}d ago"
except (ValueError, TypeError):
    freshness = "unknown"

h1, h2, h3, h4, h5 = st.columns([2, 2, 2, 2, 1])
h1.metric("Sources", f"{len(approved_sources) + 11}",
          help=f"Active scrapers: {SOURCE_LIST}. Plus {len(approved_sources)} auto-approved community sources.")
h2.metric("Signals", f"{len(insights):,}",
          help=f"{len(insights):,} quality signals from {len(_raw_insights):,} total scraped (filtered by age, relevance, and dedup).")
h3.metric("Companies", f"{len(company_meta)}")
h4.metric("Last Updated", freshness)
with h5:
    st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
    st.markdown(
        '<a href="https://github.com/grdallas-stack/geo-pulse/actions/workflows/refresh_pipeline.yml" '
        'target="_blank" style="display:inline-block; background:transparent; color:#0E3B7E; '
        'border:1px solid #0E3B7E; font-family:DM Mono,monospace; font-size:11px; '
        'text-transform:uppercase; padding:5px 12px; text-decoration:none; letter-spacing:0.03em; '
        'white-space:nowrap;" '
        'onmouseover="this.style.backgroundColor=\'#0E3B7E\';this.style.color=\'#F8F4EB\'" '
        'onmouseout="this.style.backgroundColor=\'transparent\';this.style.color=\'#0E3B7E\'"'
        '>&#8635; Refresh</a>',
        unsafe_allow_html=True,
    )

_sources_count = len(approved_sources) + 11
_provenance_ts = last_ts[:16].replace("T", " ") if last_ts else "unknown"
st.markdown(
    f'<p style="font-family:DM Mono,monospace; font-size:11px; '
    f'color:#D1CFBA; letter-spacing:0.05em;">'
    f'{len(insights):,} signals &middot; {len(company_meta)} companies &middot; '
    f'{_sources_count} sources monitored &middot; '
    f'Pipeline: {_provenance_ts}</p>',
    unsafe_allow_html=True,
)

with st.expander("New here? How to use GEO Pulse", expanded=False):
    st.markdown("""
**GEO Pulse** is a live market intelligence dashboard for the Generative Engine Optimization (GEO) and Answer Engine Optimization (AEO) category. It continuously ingests signals from practitioner forums, review sites, trade press, and Hacker News, so you can track competitor moves, buyer sentiment, and emerging feature gaps without reading hundreds of posts.

**Start with Ask AI.** Type any question and get a data-grounded intelligence brief. Examples:
- *"What are practitioners saying about Profound vs Semrush for GEO measurement?"*
- *"Which company launched something new this week?"*
- *"What features are buyers asking for most?"*

**Where the data comes from:**
""")
    _sources_df = pd.DataFrame({
        "Source": ["Hacker News", "Reddit", "G2 Reviews",
                   "Product Hunt", "Google News",
                   "Search Engine Journal / Land",
                   "Trade Press RSS"],
        "What it captures": [
            "Builder and dev discussion on GEO tools and AI search",
            "r/SEO, r/marketing, r/webdev practitioner forums",
            "Verified buyer reviews of GEO/AEO tools",
            "New tool launches in the GEO/AEO category",
            "Broad news coverage and press mentions",
            "Industry analysis and expert commentary",
            "Announcements, product updates, partnerships",
        ],
    })
    st.dataframe(_sources_df, hide_index=True, use_container_width=True)
    st.caption("Every signal is enriched with sentiment, category, and relevance scoring.")

def _data_missing():
    """Check if core data files are missing or empty."""
    if not os.path.exists(INSIGHTS_PATH):
        return True
    try:
        size = os.path.getsize(INSIGHTS_PATH)
        return size < 10  # empty json "[]" is 2 bytes
    except OSError:
        return True


def _run_pipeline_subprocess():
    """Run the pipeline as a subprocess and reload data on completion."""
    import subprocess
    import sys
    env = os.environ.copy()
    env["PYTHONPATH"] = os.path.dirname(os.path.abspath(__file__))
    result = subprocess.run(
        [sys.executable, "run_pipeline.py"],
        cwd=os.path.dirname(os.path.abspath(__file__)),
        env=env,
        capture_output=True, text=True, timeout=300,
    )
    return result.returncode == 0, result.stderr


# Auto-initialize on first load if data is missing
if _data_missing():
    if "pipeline_ran" not in st.session_state:
        with st.status("First run detected. Initializing data pipeline...", expanded=True) as status:
            st.write("Scraping Reddit, Hacker News, G2, Product Hunt, News RSS...")
            ok, err = _run_pipeline_subprocess()
            if ok:
                st.session_state.pipeline_ran = True
                status.update(label="Pipeline complete. Reloading...", state="complete")
                st.cache_data.clear()
                st.rerun()
            else:
                st.session_state.pipeline_ran = "failed"
                status.update(label="Pipeline failed.", state="error")
                st.code(err[-500:] if err else "No error output captured.")

    # Manual fallback button
    if not insights:
        st.warning("No data available. The pipeline may have failed or is still running.")
        if st.button("Initialize Data", type="primary"):
            st.session_state.pop("pipeline_ran", None)
            st.rerun()
        st.stop()
elif not insights:
    st.info("New signals ingesting. Check back in minutes.")
    st.stop()


# ---------------------------------------------------------------------------
# Ask GEO Pulse (above tabs, matching SignalSynth-GEO layout)
# ---------------------------------------------------------------------------

st.markdown("### Ask GEO Pulse")
st.caption("Ask any question about GEO/AEO market signals and get a data-grounded intelligence brief.")

from dotenv import load_dotenv
load_dotenv()
_anthropic_key = os.environ.get("ANTHROPIC_API_KEY", "")

if not _anthropic_key:
    st.warning("Anthropic API key not configured. Add `ANTHROPIC_API_KEY` to your `.env` file to enable AI Q&A.")
else:
    import re as _re

    if "qa_messages" not in st.session_state:
        st.session_state["qa_messages"] = []
    if "qa_submit" not in st.session_state:
        st.session_state["qa_submit"] = ""
    if "qa_sources_map" not in st.session_state:
        st.session_state["qa_sources_map"] = {}

    def _parse_response(raw, sources_map):
        """Split body, build clickable sources from map, extract follow-ups."""
        lines = raw.split("\n")
        body_lines = []
        followups = []
        in_sources = False
        for line in lines:
            stripped = line.strip()
            # Follow-up lines: >> question text
            if stripped.startswith(">>"):
                q = _re.sub(r'^>+\s*', '', stripped).strip()
                if q:
                    followups.append(q)
                continue
            # Detect start of Sources section and skip LLM-generated source lines
            if stripped.lower().startswith("**sources**") or stripped.lower() == "sources":
                in_sources = True
                continue
            if in_sources and _re.match(r'^\[S\d+\]', stripped):
                continue
            if in_sources and stripped == "":
                continue
            if in_sources and stripped:
                in_sources = False
            body_lines.append(line)

        body = "\n".join(body_lines).rstrip()

        # Build our own clickable sources block from the stored map
        cited = sorted(set(_re.findall(r'\[S(\d+)\]', body)), key=int)
        source_lines = []
        for num in cited:
            sid = f"S{num}"
            s = sources_map.get(sid)
            if s and s.get("url"):
                source_lines.append(f"\\[{sid}\\] [{s['title']}]({s['url']}) ({s['source']})")
            elif s:
                source_lines.append(f"\\[{sid}\\] {s['title']} ({s['source']})")

        if source_lines:
            body += "\n\n**Sources**\n\n" + "\n\n".join(source_lines)

        return body, followups

    def _do_submit(question):
        """Set the question to be processed on next rerun."""
        st.session_state["qa_submit"] = question

    # --- Input form (Enter key submits via form) ---
    with st.form("qa_form", clear_on_submit=True):
        user_question = st.text_input(
            "Ask a question",
            value="Which competitor is gaining the most momentum right now?",
            placeholder="e.g., Which competitor is gaining the most momentum right now?",
            key="ask_ai_text",
        )
        submitted = st.form_submit_button("Ask AI", type="primary")
    if submitted and user_question.strip():
        st.session_state["qa_submit"] = user_question.strip()

    # --- Pill buttons ---
    _pill_questions = [
        "What's the biggest trend this week?",
        "Where is Profound showing up?",
        "What are buyers complaining about?",
        "Which source has the most signals?",
        "What should our team know today?",
    ]

    # Render pill buttons as styled HTML, click sets query param -> triggers rerun
    _pill_html_items = ""
    for _pi, _pq in enumerate(_pill_questions):
        _encoded_q = _pq.replace("'", "&#39;").replace('"', "&quot;")
        _pill_html_items += (
            f'<a href="?pill_q={_pi}" target="_self" '
            f'style="display:inline-block; background:transparent; color:#0E3B7E; '
            f'border:1px solid #0E3B7E; font-family:DM Mono,monospace; font-size:11px; '
            f'text-transform:uppercase; padding:5px 12px; margin:0 4px 4px 0; '
            f'text-decoration:none; letter-spacing:0.03em; white-space:nowrap;"'
            f' onmouseover="this.style.backgroundColor=\'#0E3B7E\';this.style.color=\'#F8F4EB\'"'
            f' onmouseout="this.style.backgroundColor=\'transparent\';this.style.color=\'#0E3B7E\'"'
            f'>{_pq}</a>'
        )
    st.markdown(
        f'<div style="display:flex; flex-wrap:wrap; margin-top:8px;">{_pill_html_items}</div>',
        unsafe_allow_html=True,
    )

    # Check if a pill was clicked via query param
    _qp = st.query_params.get("pill_q")
    if _qp is not None:
        try:
            _pill_idx = int(_qp)
            if 0 <= _pill_idx < len(_pill_questions):
                st.session_state["qa_submit"] = _pill_questions[_pill_idx]
        except (ValueError, IndexError):
            pass
        st.query_params.clear()
        st.rerun()

    # --- Process pending question ---
    question_to_ask = st.session_state.get("qa_submit", "")
    if question_to_ask:
        st.session_state["qa_submit"] = ""
        st.session_state["qa_messages"].append({"role": "user", "content": question_to_ask})

        # Build context with source IDs
        data_summary = _build_data_summary()
        relevant_posts = _get_relevant_posts(question_to_ask)
        posts_context = ""
        sources_ref = ""
        sources_map = {}
        for idx, p in enumerate(relevant_posts[:10], 1):
            sid = f"S{idx}"
            comps = ", ".join(p.get("companies_mentioned", []))
            title = p.get("title", "")[:100] or p.get("text", "")[:60]
            url = p.get("url", "")
            source = p.get("source", "")
            posts_context += (
                f"- [{sid}] [{source}] {title} "
                f"| sentiment={p.get('sentiment','')} "
                f"| companies={comps} "
                f"| {p.get('post_date','')}\n"
                f"  \"{p.get('text','')[:200]}\"\n"
            )
            sources_ref += f"[{sid}] {title} ({source}) {url}\n"
            sources_map[sid] = {"title": title, "source": source, "url": url}

        st.session_state["qa_sources_map"] = sources_map

        system_prompt = f"""You are GEO Pulse, a market intelligence assistant for the GEO/AEO category. ProRata/Gist is the user's own product.

Your response must be boardroom-ready: concise, specific, and grounded only in provided data.
If evidence is weak, explicitly say so.

Format your answer exactly with these headings:
1) **Executive answer** (2-3 sentences, direct answer first)
2) **What the signals show** (3-6 bullets with inline citations [S1], [S2] referencing sources below)
3) **Implications for ProRata/Gist** (2-3 bullets, what this means for the product)
4) **Recommended actions** (2-4 numbered actions with owner: Product, Engineering, GTM, or Leadership)
5) **Confidence & gaps** (1-2 bullets on evidence strength and what's missing)

Do NOT include a Sources section. Sources will be rendered separately.

Use [S1], [S2] etc. inline to cite the signals below.

Rules:
- Never invent facts not present in the provided signals.
- Cite numbers: "12 mentions" not "several."
- No preamble. No "Great question." Start with the executive answer.
- No hedging, no filler, no enthusiasm. Operator-clean.
- Keep total response under ~350 words unless explicitly asked for more.
- End with 2-3 follow-up questions on new lines prefixed with ">>".

DATASET SUMMARY:
{data_summary}

RELEVANT SIGNALS:
{posts_context}

SOURCE REFERENCE (for [S1] etc. citations):
{sources_ref}
"""
        try:
            import anthropic
            client = anthropic.Anthropic(api_key=_anthropic_key)

            messages = []
            for msg in st.session_state["qa_messages"][-6:]:
                messages.append({"role": msg["role"], "content": msg["content"]})

            with st.spinner("Searching signals and generating brief..."):
                response = client.messages.create(
                    model="claude-haiku-4-5-20251001",
                    system=system_prompt,
                    messages=messages,
                    max_tokens=1200,
                    temperature=0.3,
                )

            answer = response.content[0].text
            st.session_state["qa_messages"].append({"role": "assistant", "content": answer})
        except Exception as e:
            st.session_state["qa_messages"].append({"role": "assistant", "content": f"Error: {e}"})
        st.rerun()

    # --- Render chat history ---
    if st.session_state.get("qa_messages"):
        with st.expander("AI Q&A responses", expanded=True):
            for msg in st.session_state["qa_messages"]:
                with st.chat_message(msg["role"]):
                    if msg["role"] == "assistant":
                        body, _ = _parse_response(msg["content"], st.session_state.get("qa_sources_map", {}))
                        st.markdown(body)
                    else:
                        st.markdown(msg["content"])

            # Follow-up question buttons from the last assistant message
            last_assistant = None
            for msg in reversed(st.session_state["qa_messages"]):
                if msg["role"] == "assistant":
                    last_assistant = msg["content"]
                    break
            if last_assistant:
                _, followups = _parse_response(last_assistant, {})
                if followups:
                    st.markdown("**Follow-up questions:**")
                    for fidx, fq in enumerate(followups):
                        st.button(fq, key=f"followup_{fidx}", on_click=_do_submit, args=(fq,))

            if st.button("Clear chat", key="clear_qa"):
                st.session_state["qa_messages"] = []
                st.session_state["qa_sources_map"] = {}
                st.rerun()


# ---------------------------------------------------------------------------
# Pre-compute export data at module level (no Roadmap tab dependency)
# ---------------------------------------------------------------------------

_EXPORT_OPPORTUNITY_THEMES = {
    "Real-time Tracking": ["real-time", "real time", "live tracking", "live monitoring",
                           "instant", "continuous", "live data", "monitor"],
    "Multi-LLM Coverage": ["multiple llm", "all llm", "perplexity and chatgpt", "cross-platform",
                           "every ai", "all ai", "multi-model", "chatgpt and gemini",
                           "claude and", "different ai", "llm coverage"],
    "Actionable Recs": ["actionable", "what to do", "next steps", "recommendations",
                        "how to improve", "specific advice", "optimization tip",
                        "action item"],
    "ROI Measurement": ["roi", "return on investment", "revenue impact", "attribution",
                        "prove value", "business impact", "conversion", "kpi",
                        "measure results", "performance metric"],
    "Historical Trends": ["historical", "trend", "over time", "change over", "compare week",
                          "month over month", "trajectory", "time series",
                          "tracking progress", "weekly report"],
    "Comp. Benchmarking": ["benchmark", "compare to competitor", "competitive",
                           "industry average", "how do we compare", "vs competitor",
                           "competitor analysis", "comparison", "ranking"],
    "Content Guidance": ["what to write", "content recommendations", "topic suggestion",
                         "content gap", "optimization guide", "content strategy",
                         "content plan"],
    "Brand Safety": ["brand safety", "misinformation", "hallucination about",
                     "wrong information", "incorrect", "ai says wrong",
                     "inaccurate", "false information", "reputation"],
    "Integrations": ["integrate", "integration", "connect to", "plugin",
                     "works with", "api", "google analytics", "hubspot",
                     "webhook", "zapier", "export data", "third-party"],
}

_EXPORT_FEATURE_NAMES = list(_EXPORT_OPPORTUNITY_THEMES.keys())

# Build opportunity data for exports
_export_opp_data = defaultdict(lambda: {
    "complaints": 0, "requests": 0, "praise": 0,
    "evidence": 0, "companies_tried": set(),
    "companies_praised": set(), "companies_complained": set(),
    "signals": [], "confidence": 0,
    "company_detail": defaultdict(lambda: {"count": 0, "latest": ""}),
})

for _ei in insights:
    _etxt = (_ei.get("text", "") + " " + _ei.get("title", "")).lower()
    _etags = _ei.get("entity_tags", [])
    _esent = _ei.get("sentiment", "")
    _ecomps_i = _ei.get("companies_mentioned", [])
    for _eopp, _ekws in _EXPORT_OPPORTUNITY_THEMES.items():
        if any(kw in _etxt for kw in _ekws):
            _eod = _export_opp_data[_eopp]
            _eod["evidence"] += 1
            if "complaint" in _etags:
                _eod["complaints"] += 1
                for _ec in _ecomps_i:
                    _eod["companies_complained"].add(_ec)
            if _ei.get("is_feature_request"):
                _eod["requests"] += 1
            if "praise" in _etags and _esent == "positive":
                _eod["praise"] += 1
                for _ec in _ecomps_i:
                    _eod["companies_praised"].add(_ec)
            for _ec in _ecomps_i:
                _eod["companies_tried"].add(_ec)
                _edetail = _eod["company_detail"][_ec]
                _edetail["count"] += 1
                _epd = _ei.get("post_date", "")
                if _epd > _edetail["latest"]:
                    _edetail["latest"] = _epd
            _eod["signals"].append(_ei)

for _eopp, _eod in _export_opp_data.items():
    _econf = 30
    _esrc_seen = set(s.get("source", "") for s in _eod["signals"] if s.get("source"))
    _econf += len(_esrc_seen) * 6
    _econf += min(max(_eod["evidence"] - 1, 0) * 5, 15)
    if any("G2" in s.get("source", "") for s in _eod["signals"]):
        _econf += 10
    _eod["confidence"] = min(_econf, 95)

# Build comp_stats for ALL companies at module level
_all_comp_names_export = sorted(set(
    c for i in insights for c in i.get("companies_mentioned", [])
))
_export_week_ago = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d")
_export_2week_ago = (datetime.now() - timedelta(days=14)).strftime("%Y-%m-%d")

_export_all_comp_stats = {}
for _ecomp in _all_comp_names_export:
    _esigs = [s for s in insights if _ecomp in s.get("companies_mentioned", [])]
    _etotal = len(_esigs)
    _epos = sum(1 for s in _esigs if s.get("sentiment") == "positive")
    _eneg = sum(1 for s in _esigs if s.get("sentiment") == "negative")
    _etw = sum(1 for s in _esigs if s.get("post_date", "") >= _export_week_ago)
    _epw = sum(1 for s in _esigs
               if _export_2week_ago <= s.get("post_date", "") < _export_week_ago)
    _elatest = max((s.get("post_date", "") for s in _esigs), default="")

    if _etotal < 3:
        _emom = "Insufficient data"
    elif _epw == 0 and _etw == 0:
        _emom = "No recent activity"
    elif _epw == 0:
        _emom = "Rising"
    else:
        _ewow = round((_etw - _epw) / _epw * 100)
        if _ewow >= 20:
            _emom = "Rising"
        elif _ewow <= -20:
            _emom = "Falling"
        else:
            _emom = "Stable"

    _export_all_comp_stats[_ecomp] = {
        "total": _etotal, "pos": _epos, "neg": _eneg,
        "pos_pct": round(_epos * 100 / max(_etotal, 1)),
        "momentum": _emom, "latest": _elatest,
    }


def _build_export_chart_images(all_comps, comp_stats, opp_data):
    """Build Plotly chart images for the briefing deck export."""
    momentum_colors = {
        "Rising": "#BDDEC3", "Falling": "#F44C63", "Stable": "#D1CFBA",
        "No recent activity": "#D1CFBA", "Insufficient data": "#D1CFBA",
    }

    # Filter to companies with signals
    chart_comps = [c for c in all_comps if comp_stats.get(c, {}).get("total", 0) > 0]
    sorted_comps = sorted(chart_comps, key=lambda c: comp_stats[c]["total"])

    # Fig 1: Momentum scatter
    fig1 = go.Figure()
    for m_label, m_color in momentum_colors.items():
        group = [c for c in sorted_comps if comp_stats[c]["momentum"] == m_label]
        if group:
            fig1.add_trace(go.Scatter(
                x=[comp_stats[c]["total"] for c in group],
                y=group, mode="markers",
                marker=dict(size=14, color=m_color), name=m_label,
            ))
        else:
            fig1.add_trace(go.Scatter(
                x=[None], y=[None], mode="markers",
                marker=dict(size=14, color=m_color), name=m_label, showlegend=True,
            ))
    mean_sig = sum(comp_stats[c]["total"] for c in chart_comps) / max(len(chart_comps), 1)
    fig1.add_vline(x=mean_sig, line_dash="dash", line_color="#0A0A0A", line_width=1)
    fig1.update_layout(
        plot_bgcolor="#F8F4EB", paper_bgcolor="#F8F4EB",
        font=dict(family="DM Sans, Arial, sans-serif", color="#0A0A0A"),
        xaxis=dict(title="Total signals", gridcolor="#D1CFBA", zeroline=False),
        yaxis=dict(title="", gridcolor="#D1CFBA"),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="left", x=0),
        margin=dict(l=10, r=20, t=40, b=20),
        height=max(300, len(chart_comps) * 32 + 80),
    )

    # Fig 2: Feature heatmap
    feature_names = list(_EXPORT_OPPORTUNITY_THEMES.keys())
    _30d = (datetime.now() - timedelta(days=30)).strftime("%Y-%m-%d")
    _180d = (datetime.now() - timedelta(days=180)).strftime("%Y-%m-%d")
    top_comps = chart_comps[:15] if len(chart_comps) > 15 else chart_comps

    heat_z = []
    for feat in feature_names:
        row = []
        od = opp_data.get(feat)
        for comp in top_comps:
            if od is None:
                row.append(0)
                continue
            detail = od["company_detail"].get(comp, {"count": 0, "latest": ""})
            count, latest = detail["count"], detail["latest"]
            if count == 0:
                score = 0
            elif latest and latest >= _30d:
                score = 7 + min(count - 1, 2)
            elif latest and latest >= _180d:
                score = 4 + min(count - 1, 2)
            else:
                score = 1 + min(count - 1, 2)
            row.append(min(score, 10))
        heat_z.append(row)

    fig2 = go.Figure(data=go.Heatmap(
        z=heat_z, x=top_comps, y=feature_names,
        colorscale=[[0, "#F8F4EB"], [1, "#0E3B7E"]],
        xgap=2, ygap=2, zmin=0, zmax=10,
    ))
    fig2.update_layout(
        plot_bgcolor="#F8F4EB", paper_bgcolor="#F8F4EB",
        font=dict(family="DM Sans, Arial, sans-serif", color="#0A0A0A"),
        xaxis=dict(tickangle=-35, tickfont=dict(size=11), side="bottom"),
        yaxis=dict(autorange="reversed", tickfont=dict(size=11)),
        margin=dict(l=10, r=20, t=30, b=100),
        height=max(350, len(feature_names) * 38 + 140),
    )

    try:
        img1 = fig1.to_image(format="png", width=1200, height=600, scale=2)
    except Exception:
        img1 = None
    try:
        img2 = fig2.to_image(format="png", width=1200, height=600, scale=2)
    except Exception:
        img2 = None
    return img1, img2


# ---------------------------------------------------------------------------
# Export UI (above tabs, visible on all tabs)
# ---------------------------------------------------------------------------

with st.expander("Intelligence Brief"):
    _rpt_date = datetime.now().strftime("%B %d, %Y")
    _rpt_week_ago = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d")

    # --- Pre-compute report data ---
    _rpt_week_signals = [s for s in insights if s.get("post_date", "") >= _rpt_week_ago]
    _rpt_source_counts = Counter(s.get("source", "Unknown") for s in _rpt_week_signals)
    _rpt_top_sources = _rpt_source_counts.most_common(5)
    _rpt_comp_week = Counter()
    for _rs in _rpt_week_signals:
        for _rc in _rs.get("companies_mentioned", []):
            _rpt_comp_week[_rc] += 1
    _rpt_top_comps = _rpt_comp_week.most_common(8)
    _rpt_top_signals = sorted(
        [s for s in _rpt_week_signals if _relevance_sentence(s)],
        key=lambda x: _relevance_score(x), reverse=True,
    )[:5]
    _rpt_tag_counts = Counter()
    for _rs in _rpt_week_signals:
        for _rt in _rs.get("entity_tags", []):
            _rpt_tag_counts[_rt] += 1
    _rpt_top_tags = _rpt_tag_counts.most_common(8)

    # --- Section: Title ---
    st.markdown(
        f'<div style="border-bottom:2px solid #0E3B7E; padding-bottom:12px; margin-bottom:20px;">'
        f'<h2 style="font-family:DM Sans,sans-serif; color:#0A0A0A; margin:0;">GEO Pulse Intelligence Brief</h2>'
        f'<span style="font-family:DM Mono,monospace; font-size:12px; color:#94a3b8;">{_rpt_date}</span>'
        f'</div>',
        unsafe_allow_html=True,
    )

    # --- Section 1: Market Overview ---
    _rpt_src_list = ", ".join(f"{s} ({c})" for s, c in _rpt_top_sources) if _rpt_top_sources else "N/A"
    _rpt_comp_list = ", ".join(f"{c} ({n})" for c, n in _rpt_comp_week.most_common(5)) if _rpt_comp_week else "N/A"
    st.markdown(
        f'<div style="margin-bottom:20px;">'
        f'<h3 style="font-family:DM Sans,sans-serif; color:#0E3B7E; font-size:16px; '
        f'border-bottom:1px solid #D1CFBA; padding-bottom:6px;">1. Market Overview</h3>'
        f'<table style="width:100%; font-family:DM Sans,sans-serif; font-size:14px; border-collapse:collapse;">'
        f'<tr><td style="padding:4px 0; color:#6b7280; width:180px;">Total signals this week</td>'
        f'<td style="padding:4px 0; font-weight:600;">{len(_rpt_week_signals):,}</td></tr>'
        f'<tr><td style="padding:4px 0; color:#6b7280;">Total signals all time</td>'
        f'<td style="padding:4px 0; font-weight:600;">{len(insights):,}</td></tr>'
        f'<tr><td style="padding:4px 0; color:#6b7280;">Most active sources</td>'
        f'<td style="padding:4px 0;">{_rpt_src_list}</td></tr>'
        f'<tr><td style="padding:4px 0; color:#6b7280;">Top companies this week</td>'
        f'<td style="padding:4px 0;">{_rpt_comp_list}</td></tr>'
        f'</table></div>',
        unsafe_allow_html=True,
    )

    # --- Section 2: Competitor Snapshot ---
    _rpt_comp_rows = ""
    for _rc, _rn in _rpt_top_comps:
        _rcs = _export_all_comp_stats.get(_rc, {})
        _rcmom = _rcs.get("momentum", "Unknown")
        _rcpos = _rcs.get("pos_pct", 0)
        _rpt_comp_rows += (
            f'<tr><td style="padding:5px 0; font-weight:600;">{_rc}</td>'
            f'<td style="padding:5px 0; text-align:center;">{_rn}</td>'
            f'<td style="padding:5px 0; text-align:center;">{_rcmom}</td>'
            f'<td style="padding:5px 0; text-align:center;">{_rcpos}% positive</td></tr>'
        )
    st.markdown(
        f'<div style="margin-bottom:20px;">'
        f'<h3 style="font-family:DM Sans,sans-serif; color:#0E3B7E; font-size:16px; '
        f'border-bottom:1px solid #D1CFBA; padding-bottom:6px;">2. Competitor Snapshot</h3>'
        f'<table style="width:100%; font-family:DM Sans,sans-serif; font-size:13px; border-collapse:collapse;">'
        f'<tr style="color:#94a3b8; font-family:DM Mono,monospace; font-size:11px; text-transform:uppercase;">'
        f'<th style="text-align:left; padding:6px 0; border-bottom:1px solid #D1CFBA;">Company</th>'
        f'<th style="text-align:center; padding:6px 0; border-bottom:1px solid #D1CFBA;">Signals</th>'
        f'<th style="text-align:center; padding:6px 0; border-bottom:1px solid #D1CFBA;">Momentum</th>'
        f'<th style="text-align:center; padding:6px 0; border-bottom:1px solid #D1CFBA;">Sentiment</th></tr>'
        f'{_rpt_comp_rows}</table></div>',
        unsafe_allow_html=True,
    )

    # --- Section 3: Top Signals This Week ---
    _rpt_sig_html = ""
    for _rsi, _rs in enumerate(_rpt_top_signals, 1):
        _rst = (_rs.get("title", "") or _rs.get("text", ""))[:120]
        _rsu = _rs.get("url", "")
        _rss = _rs.get("source", "Unknown")
        _rsd = _time_ago(_rs.get("post_date", ""))
        _rsb = _relevance_sentence(_rs) or ""
        _rst_link = f'<a href="{_rsu}" target="_blank" style="color:#0A0A0A; text-decoration:none; font-weight:600;">{_rst}</a>' if _rsu else f'<b>{_rst}</b>'
        _rpt_sig_html += (
            f'<div style="padding:10px 0; border-bottom:1px solid #E8E4D9;">'
            f'<div style="display:flex; justify-content:space-between;">'
            f'<span style="font-family:DM Mono,monospace; font-size:10px; color:#94a3b8; text-transform:uppercase;">{_rss}</span>'
            f'<span style="font-family:DM Mono,monospace; font-size:10px; color:#94a3b8;">{_rsd}</span></div>'
            f'<div style="margin-top:4px; font-size:14px;">{_rst_link}</div>'
            f'<div style="margin-top:2px; font-size:12px; color:#6b7280;">{_rsb}</div>'
            f'</div>'
        )
    st.markdown(
        f'<div style="margin-bottom:20px;">'
        f'<h3 style="font-family:DM Sans,sans-serif; color:#0E3B7E; font-size:16px; '
        f'border-bottom:1px solid #D1CFBA; padding-bottom:6px;">3. Top Signals This Week</h3>'
        f'{_rpt_sig_html}</div>',
        unsafe_allow_html=True,
    )

    # --- Section 4: Trends & Themes ---
    _rpt_tag_html = ""
    _tag_labels = {
        "buyer_voice": "Buyer discussions driving demand signals",
        "founder_voice": "Founders sharing product and GTM insights",
        "analyst_voice": "Analyst and media coverage shaping narrative",
        "feature_request": "Feature gaps surfacing in buyer conversations",
        "competitive_intel": "Competitive positioning and comparison activity",
        "complaint": "Pain points and dissatisfaction signals",
        "praise": "Positive reception and advocacy signals",
        "funding_news": "Capital activity signaling market confidence",
        "product_launch": "New products and feature releases entering market",
    }
    for _rtt, _rtc in _rpt_top_tags:
        _rtl = _tag_labels.get(_rtt, _rtt.replace("_", " ").title())
        _rpt_tag_html += (
            f'<tr><td style="padding:4px 0;">'
            f'<span style="display:inline-block; background:#0E3B7E; color:#F8F4EB; '
            f'font-family:DM Mono,monospace; font-size:10px; padding:2px 8px; '
            f'text-transform:uppercase; letter-spacing:0.03em;">{_rtt.replace("_"," ")}</span></td>'
            f'<td style="padding:4px 8px; text-align:center; font-weight:600;">{_rtc}</td>'
            f'<td style="padding:4px 0; color:#6b7280; font-size:13px;">{_rtl}</td></tr>'
        )
    st.markdown(
        f'<div style="margin-bottom:20px;">'
        f'<h3 style="font-family:DM Sans,sans-serif; color:#0E3B7E; font-size:16px; '
        f'border-bottom:1px solid #D1CFBA; padding-bottom:6px;">4. Trends &amp; Themes</h3>'
        f'<table style="width:100%; font-family:DM Sans,sans-serif; font-size:14px; border-collapse:collapse;">'
        f'{_rpt_tag_html}</table></div>',
        unsafe_allow_html=True,
    )

    # --- Section 5: Recommended Actions ---
    _rpt_actions = []
    if _rpt_top_tags and _rpt_top_tags[0][1] >= 3:
        _top_theme = _rpt_top_tags[0][0].replace("_", " ")
        _rpt_actions.append(
            f"Investigate the surge in <b>{_top_theme}</b> signals ({_rpt_top_tags[0][1]} this week). "
            f"Determine if this represents a positioning opportunity or emerging threat."
        )
    if _rpt_top_comps:
        _top_c = _rpt_top_comps[0][0]
        _top_cn = _rpt_top_comps[0][1]
        _rpt_actions.append(
            f"<b>{_top_c}</b> leads competitor mentions this week ({_top_cn} signals). "
            f"Review their recent activity for positioning shifts or product launches that warrant a response."
        )
    _neg_comps = [c for c, s in _export_all_comp_stats.items()
                  if s.get("pos_pct", 100) < 40 and s.get("total", 0) >= 3 and _rpt_comp_week.get(c, 0) >= 2]
    if _neg_comps:
        _rpt_actions.append(
            f"Competitors with low positive sentiment ({', '.join(_neg_comps[:3])}) may have exploitable gaps. "
            f"Cross-reference with buyer complaints to identify differentiation angles."
        )
    if len(_rpt_actions) < 3:
        _rpt_actions.append(
            "Review the top 5 signals above for messaging or content opportunities. "
            "Prioritize any that mention ProRata/Gist directly or reference unmet buyer needs."
        )

    _rpt_action_html = "".join(
        f'<li style="margin-bottom:8px; font-size:14px; line-height:1.5;">{a}</li>'
        for a in _rpt_actions[:3]
    )
    st.markdown(
        f'<div style="margin-bottom:20px;">'
        f'<h3 style="font-family:DM Sans,sans-serif; color:#0E3B7E; font-size:16px; '
        f'border-bottom:1px solid #D1CFBA; padding-bottom:6px;">5. Recommended Actions</h3>'
        f'<ol style="font-family:DM Sans,sans-serif; padding-left:20px;">{_rpt_action_html}</ol>'
        f'</div>',
        unsafe_allow_html=True,
    )

    # --- Download as .docx ---
    st.markdown("---")
    _date_tag = datetime.now().strftime("%Y-%m-%d")
    try:
        _buf = _export_research_report(
            insights, company_meta, _export_opp_data,
            _all_comp_names_export, _export_all_comp_stats
        )
        _doc_bytes = _buf.getvalue()
        st.download_button(
            label="Download Report (.docx)",
            data=_doc_bytes,
            file_name=f"GEOPulse_IntelligenceBrief_{_date_tag}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            key="export_dl_brief",
            type="primary",
        )
    except Exception as _ex:
        st.error(f"Export failed: {_ex}")


# ---------------------------------------------------------------------------
# URL state restoration from query params
# ---------------------------------------------------------------------------
_qp = st.query_params
_shared_section = _qp.get("section", "")
_shared_tab_map = {
    "live_feed": 0, "signal_of_week": 0,
    "competitors": 1,
    "momentum_chart": 2, "heat_map": 2, "build_now": 2,
}
_shared_tab_idx = _shared_tab_map.get(_shared_section, None)

# Restore competitor filter from URL if present
if "comps" in _qp:
    _url_comps = _qp.get("comps", "").split(",")
    if _url_comps and _url_comps[0]:
        st.session_state["rm_comp_sel"] = _url_comps

# ---------------------------------------------------------------------------
# Tabs (3 tabs)
# ---------------------------------------------------------------------------
_tab_names = ["Live Feed", "Competitors", "Roadmap"]
if _shared_tab_idx is not None:
    tabs = st.tabs(_tab_names)
else:
    tabs = st.tabs(_tab_names)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TAB 1: LIVE FEED
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
with tabs[0]:
    st.markdown("### Live Feed")
    st.markdown(
        "Real-time stream from Reddit, Hacker News, G2, Product Hunt, and trade press. "
        "Ranked by relevance, recency, and social signal strength."
    )

    # --- Signal of the Week ---
    # Dynamically pick highest-relevance signal from last 48 hours, fall back to editorial pick
    _sotw_cutoff = (datetime.now() - timedelta(hours=48)).strftime("%Y-%m-%d")
    _sotw_candidates = [
        s for s in insights
        if s.get("post_date", "") >= _sotw_cutoff
        and _is_display_relevant(s)
        and _is_displayable_post(s)
        and (s.get("companies_mentioned") or s.get("entity_tags") or s.get("is_competitive_intel"))
        and _relevance_score(s) >= _MIN_RELEVANCE_SCORE * 2
    ]
    _sotw_candidates.sort(key=lambda x: _relevance_score(x), reverse=True)

    if _sotw_candidates:
        _sotw_pick = _sotw_candidates[0]
        _sotw_company = (_sotw_pick.get("companies_mentioned") or [""])[0]
        _sotw_title = (_sotw_pick.get("title") or _sotw_pick.get("text", ""))[:150]
        _sotw_url = _sotw_pick.get("url", "")
        _sotw_brief = (_sotw_pick.get("signal_brief") or "").strip() or "High-relevance signal for GEO/AEO practitioners."
        _sotw_source = _sotw_pick.get("source", "")
        _sotw_date = _sotw_pick.get("post_date", "")
    else:
        # Fallback to editorial pick
        _sotw_company = _SOTW.get("company", "")
        _sotw_title = _SOTW.get("title", "Signal of the Week")
        _sotw_url = _SOTW.get("url", "")
        _sotw_brief = _SOTW.get("brief", "")
        _sotw_source = _SOTW.get("source", "")
        _sotw_date = _SOTW.get("date", "")

    # Ensure title is always clickable
    if _sotw_url:
        _sotw_title_html = (
            f'<a href="{_sotw_url}" target="_blank" style="font-size: 1.05rem; '
            f'font-weight: 600; color: #0A0A0A; text-decoration: none;">{_sotw_title}</a>'
        )
    else:
        _sotw_title_html = f'<span style="font-size: 1.05rem; font-weight: 600; color: #0A0A0A;">{_sotw_title}</span>'

    _sotw_comp_pill = ""
    if _sotw_company:
        _sotw_comp_pill = (
            f'<span style="display:inline-block; background:#0E3B7E; color:#F8F4EB; '
            f"font-family:'DM Mono',monospace; font-size:11px; padding:2px 8px; "
            f'text-transform:uppercase; letter-spacing:0.03em;">{_sotw_company.upper()}</span><br>'
        )

    _sotw_source_line = ""
    if _sotw_source or _sotw_date:
        _src_link = f'<a href="{_sotw_url}" target="_blank" style="color: #D1CFBA;">{_sotw_source}</a>' if _sotw_url and _sotw_source else _sotw_source
        _sotw_source_line = f'<span style="font-size: 0.75rem; color: #D1CFBA;">{_src_link} &middot; {_sotw_date}</span>'

    st.markdown(
        f'<div style="border-left: 4px solid #FF9D1C; padding: 16px; '
        f'background: #FFFFFF; border: 1px solid #D1CFBA; margin-bottom: 1rem;">'
        f'<span style="color: #FF9D1C; font-family: DM Mono, monospace; font-size: 0.75rem; font-weight: 600; '
        f'letter-spacing: 0.05em; text-transform: uppercase;">'
        f'SIGNAL OF THE WEEK</span><br>'
        f'{_sotw_comp_pill}'
        f'{_sotw_title_html}'
        f'<p style="margin: 0.5rem 0 0.4rem 0; font-size: 0.9rem; color: #0A0A0A;">'
        f'{_sotw_brief}</p>'
        f'{_sotw_source_line}'
        f'</div>',
        unsafe_allow_html=True,
    )
    # Share button inside the card area, styled as ghost pill
    _share_button("signal_of_week", "Share Signal of the Week")

    _sotw_insight = {
        "source": _sotw_source,
        "title": _sotw_title,
        "url": _sotw_url,
        "post_date": _sotw_date,
        "companies_mentioned": [_sotw_company] if _sotw_company else [],
        "entity_tags": ["product_launch"],
        "sentiment": "positive",
    }

    fc1, fc2 = st.columns(2)
    all_companies_in_data = sorted(set(
        c for i in insights for c in i.get("companies_mentioned", [])
    ))
    with fc1:
        filter_company = st.selectbox("Company", ["All"] + all_companies_in_data, key="feed_company")
    with fc2:
        sources_in_data = sorted(set(i.get("source", "") for i in insights))
        filter_source = st.selectbox("Source", ["All"] + sources_in_data, key="feed_source")

    with st.expander("Add a New Competitor to Monitor"):
        new_comp_name = st.text_input("Company name", key="new_comp_name")
        new_comp_url = st.text_input(
            "Company website URL (optional but recommended)", key="new_comp_url",
            help="Helps us find the right company, especially for common names.",
        )
        new_comp_aliases = st.text_input("Aliases (comma-separated)", key="new_comp_aliases",
                                         help="e.g. 'acme, acme.ai, acme seo'")
        if st.button("Start Monitoring", key="add_comp_btn"):
            if new_comp_name.strip():
                name = new_comp_name.strip()
                aliases = [a.strip().lower() for a in new_comp_aliases.split(",") if a.strip()]
                if not aliases:
                    aliases = [name.lower()]
                url = new_comp_url.strip()
                new_entry = {"name": name, "aliases": aliases, "category": "unknown", "url": url}
                cd = load_companies()
                cd.setdefault("competitors", []).append(new_entry)
                with open(COMPANIES_PATH, "w") as f:
                    json.dump(cd, f, ensure_ascii=False, indent=2)
                st.cache_data.clear()
                st.success(f"Now monitoring **{name}**. Will appear in next pipeline run.")

    with st.expander("Daily Email Digest"):
        st.markdown("Get GEO Pulse insights delivered to your inbox daily.")

        _subs = _load_subscribers()
        _sub_emails = {s["email"] for s in _subs}

        _d_col1, _d_col2 = st.columns(2)
        with _d_col1:
            _sub_email = st.text_input("Email address", key="sub_email")
            _sub_name = st.text_input("Name (optional)", key="sub_name")
        with _d_col2:
            _sub_hour = st.selectbox(
                "Delivery time",
                options=list(range(5, 22)),
                index=3,  # 8 AM default
                format_func=lambda h: f"{h}:00",
                key="sub_hour",
            )
            _tz_options = {
                "US/Eastern (UTC-5)": -5, "US/Central (UTC-6)": -6,
                "US/Mountain (UTC-7)": -7, "US/Pacific (UTC-8)": -8,
                "UTC": 0, "Europe/London (UTC+0)": 0,
                "Europe/Berlin (UTC+1)": 1, "Asia/Tokyo (UTC+9)": 9,
            }
            _sub_tz_label = st.selectbox("Timezone", options=list(_tz_options.keys()), key="sub_tz")
            _sub_tz_offset = _tz_options[_sub_tz_label]

        _all_comp_digest = sorted(set(
            c for i in insights for c in i.get("companies_mentioned", [])
        ))
        _sub_comp_filter = st.multiselect(
            "Competitor filter (leave empty for all)",
            options=_all_comp_digest,
            key="sub_comp_filter",
        )

        _sub_col1, _sub_col2 = st.columns(2)
        with _sub_col1:
            if st.button("Subscribe", key="subscribe_btn", type="primary"):
                if _sub_email and "@" in _sub_email:
                    if _sub_email in _sub_emails:
                        st.warning("This email is already subscribed.")
                    else:
                        new_sub = {
                            "email": _sub_email,
                            "name": _sub_name,
                            "delivery_hour": _sub_hour,
                            "tz_offset": _sub_tz_offset,
                            "competitor_filter": _sub_comp_filter,
                            "confirmed": True,
                            "subscribed_at": datetime.now().isoformat(),
                        }
                        _subs.append(new_sub)
                        _save_subscribers(_subs)
                        _send_confirmation_email(_sub_email, _sub_name)
                        st.success(f"Subscribed {_sub_email}. Confirmation email sent.")
                else:
                    st.error("Enter a valid email address.")

        with _sub_col2:
            if st.button("Unsubscribe", key="unsub_btn"):
                if _sub_email and _sub_email in _sub_emails:
                    _subs = [s for s in _subs if s["email"] != _sub_email]
                    _save_subscribers(_subs)
                    st.success(f"Unsubscribed {_sub_email}.")
                elif _sub_email:
                    st.info("Email not found in subscriber list.")

        if _subs:
            st.caption(f"{len(_subs)} active subscriber(s)")

    filtered = insights
    if filter_company != "All":
        filtered = [i for i in filtered if filter_company in i.get("companies_mentioned", [])]
    if filter_source != "All":
        filtered = [i for i in filtered if i.get("source", "") == filter_source]

    filtered = [i for i in filtered
                if _is_display_relevant(i) and (i.get("signal_brief") or "").strip()
                and _relevance_score(i) >= _MIN_RELEVANCE_SCORE]
    filtered.sort(key=lambda x: (x.get("post_date", ""), _relevance_score(x)), reverse=True)
    st.caption(f"Showing {min(25, len(filtered))} of {len(filtered)} GEO-relevant signals from {len(insights):,} total ingested (filtered for relevance)")

    new_companies = _get_new_companies(json.dumps(insights))
    _base_page_size = 25
    if "feed_page_count" not in st.session_state:
        st.session_state["feed_page_count"] = 1
    page_size = _base_page_size * st.session_state["feed_page_count"]
    for idx, insight in enumerate(filtered[:page_size]):
        source = insight.get("source", "")
        source_label = _source_badge(source)
        title = insight.get("title", "")[:120] or insight.get("text", "")[:120]
        companies = insight.get("companies_mentioned", [])
        url = insight.get("url", "")
        date = insight.get("post_date", "")
        time_label = _time_ago(date)
        rel_sentence = (insight.get("signal_brief") or "").strip()

        # Title as HTML anchor or plain text
        if url:
            title_html = (
                f'<a href="{url}" target="_blank" '
                f'style="color:#0A0A0A; text-decoration:none;">{title}</a>'
            )
        else:
            title_html = title

        # Keyword callout pills
        kw_pills = ""
        kws = _keywords_for_card(insight)
        if kws:
            kw_tags = "".join(
                f'<span style="display:inline-block; background:transparent; color:#0E3B7E; '
                f'border:1px solid #0E3B7E; font-family:DM Mono,monospace; font-size:10px; '
                f'padding:2px 8px; margin-right:4px; margin-top:4px; letter-spacing:0.03em; '
                f'text-transform:uppercase;">{kw}</span>'
                for kw in kws[:3]
            )
            kw_pills = f'<div style="margin-top:6px;">{kw_tags}</div>'

        # Competitor pills — show ALL competitors mentioned
        comp_pills = ""
        if companies:
            pills = "".join(
                f'<span style="display:inline-block; background:#0E3B7E; color:#F8F4EB; '
                f'font-family:DM Mono,monospace; font-size:10px; padding:2px 8px; '
                f'margin-right:4px; letter-spacing:0.03em;">{c}</span>'
                for c in companies
            )
            comp_pills = (
                f'<div style="margin-top:6px;">'
                f'<span style="font-family:DM Mono,monospace; font-size:10px; '
                f'color:#94a3b8; margin-right:4px;">Competitors:</span>{pills}</div>'
            )

        st.markdown(
            f'<div style="background:#FFFFFF; padding:16px; '
            f'border-bottom:1px solid #D1CFBA;">'
            f'<div style="display:flex; justify-content:space-between; align-items:center;">'
            f'<span style="font-family:DM Mono,monospace; font-size:10px; '
            f'color:#94a3b8; text-transform:uppercase; letter-spacing:0.05em;">'
            f'{source_label}</span>'
            f'<span style="font-family:DM Mono,monospace; font-size:10px; '
            f'color:#94a3b8;">{time_label}</span>'
            f'</div>'
            f'<div style="margin-top:6px; font-family:DM Sans,sans-serif; '
            f'font-size:15px; font-weight:700; color:#0A0A0A; line-height:1.3;">'
            f'{title_html}</div>'
            f'<div style="margin-top:4px; font-size:13px; color:#6b7280; '
            f'line-height:1.4;">{rel_sentence}</div>'
            f'{kw_pills}'
            f'{comp_pills}'
            f'</div>',
            unsafe_allow_html=True,
        )

    if len(filtered) > page_size:
        _remaining = len(filtered) - page_size
        _next_batch = min(25, _remaining)
        st.markdown(
            f'<div style="text-align:center; margin-top:12px;">',
            unsafe_allow_html=True,
        )
        if st.button(
            f"Load {_next_batch} more signals \u2193",
            key="load_more_signals",
        ):
            st.session_state["feed_page_count"] += 1
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
        st.caption(f"{_remaining} more signals available")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TAB 2: COMPETITORS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
with tabs[1]:
    st.markdown("### Competitors")
    _share_button("competitors", "Share Competitors")
    st.markdown(
        "One card per competitor. Momentum, sentiment, and the most relevant signals this week."
    )

    # Build per-company data
    now = datetime.now()
    week_ago_str = (now - timedelta(days=7)).strftime("%Y-%m-%d")

    comp_data_map = defaultdict(lambda: {
        "total": 0, "this_week": 0,
        "pos": 0, "neg": 0, "neu": 0,
        "signals": [],  # all signals for this company, scored
    })

    for i in insights:
        for comp in i.get("companies_mentioned", []):
            cd = comp_data_map[comp]
            cd["total"] += 1
            s = i.get("sentiment", "neutral")
            cd[{"positive": "pos", "negative": "neg"}.get(s, "neu")] += 1
            date = i.get("post_date", "")
            if date >= week_ago_str:
                cd["this_week"] += 1
            cd["signals"].append(i)

    # Sort by total mentions
    sorted_comps = sorted(comp_data_map.items(), key=lambda x: x[1]["total"], reverse=True)
    two_weeks_ago_str = (now - timedelta(days=14)).strftime("%Y-%m-%d")

    # Pre-compute 30-day window for activity checks
    _30d_ago_comp = (now - timedelta(days=30)).strftime("%Y-%m-%d")

    # Rank and sort: separate high-signal from low-signal competitors
    _ranked_comps = []
    for comp, cd in sorted_comps:
        total = cd["total"]
        if total < 1:
            continue
        _ranked_comps.append((comp, cd))

    for _rank_idx, (comp, cd) in enumerate(_ranked_comps):
        total = cd["total"]
        pos_pct = round(cd["pos"] * 100 / max(total, 1))
        neg_pct = round(cd["neg"] * 100 / max(total, 1))
        neu_pct = 100 - pos_pct - neg_pct

        # Momentum: compute WoW from actual signal dates
        this_week = cd["this_week"]
        prior_week = sum(
            1 for s in cd["signals"]
            if two_weeks_ago_str <= s.get("post_date", "") < week_ago_str
        )

        if total < 3:
            momentum = "Insufficient data"
        elif prior_week == 0 and this_week == 0:
            momentum = "No recent activity"
        elif prior_week == 0:
            momentum = f"Rising (+{this_week} new)"
        else:
            wow_pct = round((this_week - prior_week) / prior_week * 100)
            if wow_pct >= 20:
                momentum = f"Rising ({wow_pct:+}% WoW)"
            elif wow_pct <= -20:
                momentum = f"Falling ({wow_pct:+}% WoW)"
            else:
                momentum = f"Stable ({wow_pct:+}% WoW)"

        is_own = comp in own_brands
        meta = company_meta.get(comp, {})
        positioning = _company_positioning(meta) if meta else ""
        own_tag = " own brand" if is_own else ""

        # Find most recent post date for this company
        latest_date = ""
        for s in cd["signals"]:
            d = s.get("post_date", "")
            if d > latest_date:
                latest_date = d

        # Item 9 & 10: Relative time for last signal
        _last_signal_label = ""
        if latest_date:
            _last_signal_ago = _time_ago(latest_date)
            if _last_signal_ago:
                _last_signal_label = f"Last signal: {_last_signal_ago}"
            else:
                _last_signal_label = f"Last signal: {latest_date}"
        # Only show "No recent activity" if zero signals in 30 days
        _signals_30d = sum(1 for s in cd["signals"] if s.get("post_date", "") >= _30d_ago_comp)
        if _signals_30d == 0 and total > 0:
            _last_signal_label = "No recent activity"

        # Item 12: Visual hierarchy - de-emphasize low-signal cards
        _is_low_signal = total < 3
        _border_style = "1px solid #E8E4D9" if _is_low_signal else "1px solid #D1CFBA"

        # Item 12: Rank badge for top 3
        _rank_badge = ""
        if _rank_idx < 3:
            _rank_badge = f"#{_rank_idx + 1} "

        # Item 11: Signal counts instead of broken sentiment
        _all_neutral = (pos_pct == 0 and neg_pct == 0) or neu_pct == 100
        _signals_this_week = cd["this_week"]
        _signals_30d_count = _signals_30d

        with st.container(border=True):
            # Header row
            hc1, hc2 = st.columns([3, 2])
            with hc1:
                comp_url = meta.get("url", "") if meta else ""
                site_link = f" [Visit site]({comp_url})" if comp_url else ""
                _header_size = "**" if not _is_low_signal else ""
                st.markdown(f"{_rank_badge}{_header_size}{comp}{_header_size}{own_tag}{site_link}")
                if positioning:
                    st.caption(positioning)
            with hc2:
                st.markdown(f"{total} mentions, {momentum}")
                # Item 11: Hide sentiment if all neutral, show signal counts instead
                if _all_neutral:
                    st.caption(f"{_signals_this_week} signals this week, {_signals_30d_count} signals last 30 days")
                else:
                    st.caption(f"{pos_pct}% positive, {neg_pct}% negative")
                    st.caption(f"{_signals_this_week} signals this week, {_signals_30d_count} signals last 30 days")

            # Item 12: Replace "Limited data" disclaimer
            if _is_low_signal:
                st.markdown(
                    f'<span style="font-family:DM Mono,monospace; font-size:10px; '
                    f'color:#94a3b8;">Fewer than 3 signals tracked</span>',
                    unsafe_allow_html=True,
                )

            # Item 9 & 10: Show relative time label
            if _last_signal_label:
                st.caption(_last_signal_label)

            # Top 3 most relevant signals this week
            week_signals = [s for s in cd["signals"] if s.get("post_date", "") >= week_ago_str]
            if not week_signals:
                week_signals = cd["signals"]
            week_signals.sort(key=lambda x: _relevance_score(x), reverse=True)

            top_signals = [s for s in week_signals if _relevance_sentence(s, for_company=comp)][:3]
            for sig in top_signals:
                sig_title = sig.get("title", "")[:100] or sig.get("text", "")[:100]
                sig_url = sig.get("url", "")
                sig_source = sig.get("source", "")
                sig_source_label = _source_badge(sig_source)
                sig_reason = _relevance_sentence(sig, for_company=comp)

                headline = f"[{sig_title}]({sig_url})" if sig_url else sig_title
                st.markdown(f"  `{sig_source_label}` {headline}")
                st.caption(f"  _{sig_reason}_ , {_time_ago(sig.get('post_date', ''))}")

            # Show more expander
            remaining = [s for s in week_signals[3:] if _relevance_sentence(s, for_company=comp)][:12]
            if remaining:
                with st.expander(f"Show {len(remaining)} more signals"):
                    for sig in remaining:
                        sig_title = sig.get("title", "")[:100] or sig.get("text", "")[:100]
                        sig_url = sig.get("url", "")
                        sig_source = sig.get("source", "")
                        sig_source_label = _source_badge(sig_source)
                        sig_reason = _relevance_sentence(sig, for_company=comp)

                        headline = f"[{sig_title}]({sig_url})" if sig_url else sig_title
                        st.markdown(f"`{sig_source_label}` {headline}")
                        st.caption(f"_{sig_reason}_ , {_time_ago(sig.get('post_date', ''))}")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TAB 3: ROADMAP
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
with tabs[2]:
    st.markdown("### Roadmap")
    st.markdown(
        "Product opportunities ranked by market evidence, trending features, "
        "and competitive coverage."
    )

    OPPORTUNITY_THEMES = {
        "Real-time Tracking": ["real-time", "real time", "live tracking", "live monitoring",
                               "instant", "continuous", "live data", "monitor"],
        "Multi-LLM Coverage": ["multiple llm", "all llm", "perplexity and chatgpt", "cross-platform",
                               "every ai", "all ai", "multi-model", "chatgpt and gemini",
                               "claude and", "different ai", "llm coverage"],
        "Actionable Recs": ["actionable", "what to do", "next steps", "recommendations",
                            "how to improve", "specific advice", "optimization tip",
                            "action item"],
        "ROI Measurement": ["roi", "return on investment", "revenue impact", "attribution",
                            "prove value", "business impact", "conversion", "kpi",
                            "measure results", "performance metric"],
        "Historical Trends": ["historical", "trend", "over time", "change over", "compare week",
                              "month over month", "trajectory", "time series",
                              "tracking progress", "weekly report"],
        "Comp. Benchmarking": ["benchmark", "compare to competitor", "competitive",
                               "industry average", "how do we compare", "vs competitor",
                               "competitor analysis", "comparison", "ranking"],
        "Content Guidance": ["what to write", "content recommendations", "topic suggestion",
                             "content gap", "optimization guide", "content strategy",
                             "content plan"],
        "Brand Safety": ["brand safety", "misinformation", "hallucination about",
                         "wrong information", "incorrect", "ai says wrong",
                         "inaccurate", "false information", "reputation"],
        "Integrations": ["integrate", "integration", "connect to", "plugin",
                         "works with", "api", "google analytics", "hubspot",
                         "webhook", "zapier", "export data", "third-party"],
    }

    # Build opportunity data with per-company tracking and full signal refs
    opportunity_data = defaultdict(lambda: {
        "complaints": 0, "requests": 0, "praise": 0,
        "evidence": 0, "companies_tried": set(),
        "companies_praised": set(), "companies_complained": set(),
        "signals": [], "confidence": 0,
        "company_detail": defaultdict(lambda: {"count": 0, "latest": ""}),
    })

    for i in insights:
        text_lower = (i.get("text", "") + " " + i.get("title", "")).lower()
        tags = i.get("entity_tags", [])
        sentiment = i.get("sentiment", "")
        companies_i = i.get("companies_mentioned", [])

        for opp, keywords in OPPORTUNITY_THEMES.items():
            if any(kw in text_lower for kw in keywords):
                od = opportunity_data[opp]
                od["evidence"] += 1
                if "complaint" in tags:
                    od["complaints"] += 1
                    for c in companies_i:
                        od["companies_complained"].add(c)
                if i.get("is_feature_request"):
                    od["requests"] += 1
                if "praise" in tags and sentiment == "positive":
                    od["praise"] += 1
                    for c in companies_i:
                        od["companies_praised"].add(c)
                for c in companies_i:
                    od["companies_tried"].add(c)
                    detail = od["company_detail"][c]
                    detail["count"] += 1
                    post_date = i.get("post_date", "")
                    if post_date > detail["latest"]:
                        detail["latest"] = post_date
                od["signals"].append(i)

    for opp, od in opportunity_data.items():
        conf = 30
        sources_seen = set(s.get("source", "") for s in od["signals"] if s.get("source"))
        conf += len(sources_seen) * 6
        extra_signals = max(od["evidence"] - 1, 0)
        conf += min(extra_signals * 5, 15)
        if any("G2" in s.get("source", "") for s in od["signals"]):
            conf += 10
        od["confidence"] = min(conf, 95)

    # --- Shared computation ---
    comp_mention_counts = Counter()
    for i in insights:
        for c in i.get("companies_mentioned", []):
            comp_mention_counts[c] += 1
    all_companies_ranked = [c for c, _ in comp_mention_counts.most_common()]
    top8_default = [c for c, _ in comp_mention_counts.most_common(8)]

    # Date boundaries
    _now_rm = datetime.now()
    _30d_ago = (_now_rm - timedelta(days=30)).strftime("%Y-%m-%d")
    _90d_ago = (_now_rm - timedelta(days=90)).strftime("%Y-%m-%d")
    _week_ago_rm = (_now_rm - timedelta(days=7)).strftime("%Y-%m-%d")
    _2week_ago_rm = (_now_rm - timedelta(days=14)).strftime("%Y-%m-%d")

    # ── COMPETITOR SELECTOR (Item 14: styled filter buttons) ──────────────────
    if "rm_quick" not in st.session_state:
        st.session_state["rm_quick"] = None

    def _set_quick(mode):
        st.session_state["rm_quick"] = mode
        st.session_state.pop("rm_comp_sel", None)

    # Compute rising companies
    _rising_comps = []
    for c in all_companies_ranked:
        tw = sum(1 for s in insights
                 if c in s.get("companies_mentioned", [])
                 and s.get("post_date", "") >= _week_ago_rm)
        pw = sum(1 for s in insights
                 if c in s.get("companies_mentioned", [])
                 and _2week_ago_rm <= s.get("post_date", "") < _week_ago_rm)
        if pw > 0 and tw > pw:
            _rising_comps.append(c)
        elif pw == 0 and tw > 0:
            _rising_comps.append(c)

    qmode = st.session_state.get("rm_quick")

    # Item 14: Filter buttons
    qc1, qc2, qc3 = st.columns(3)
    with qc1:
        st.button("Top 8 by volume", key="qf_top8", on_click=_set_quick, args=("top8",))
    with qc2:
        st.button("Rising only", key="qf_rising", on_click=_set_quick, args=("rising",))
    with qc3:
        st.button("All", key="qf_all", on_click=_set_quick, args=("all",))

    if qmode == "top8":
        _sel_default = top8_default
    elif qmode == "rising":
        _sel_default = _rising_comps if _rising_comps else top8_default
    elif qmode == "all":
        _sel_default = all_companies_ranked
    else:
        _sel_default = top8_default

    if qmode:
        st.session_state["rm_quick"] = None

    selected_comps = st.multiselect(
        "Select competitors to display",
        options=all_companies_ranked,
        default=_sel_default,
        help="Choose which competitors appear in the tables below",
        key="rm_comp_sel",
    )

    if len(selected_comps) < 2:
        st.warning("Select at least 2 competitors to compare.")
        st.stop()

    # ── Pre-compute per-company stats ────────────────────────────────────────
    _comp_stats = {}
    for comp in selected_comps:
        sigs = [s for s in insights if comp in s.get("companies_mentioned", [])]
        total = len(sigs)
        pos = sum(1 for s in sigs if s.get("sentiment") == "positive")
        neg = sum(1 for s in sigs if s.get("sentiment") == "negative")
        tw = sum(1 for s in sigs if s.get("post_date", "") >= _week_ago_rm)
        pw = sum(1 for s in sigs
                 if _2week_ago_rm <= s.get("post_date", "") < _week_ago_rm)
        latest = max((s.get("post_date", "") for s in sigs), default="")

        if total < 3:
            momentum = "Insufficient data"
        elif pw == 0 and tw == 0:
            momentum = "No recent activity"
        elif pw == 0:
            momentum = "Rising"
        else:
            wow = round((tw - pw) / pw * 100)
            if wow >= 20:
                momentum = "Rising"
            elif wow <= -20:
                momentum = "Falling"
            else:
                momentum = "Stable"

        _30d_count = sum(1 for s in sigs if s.get("post_date", "") >= _30d_ago)

        _comp_stats[comp] = {
            "total": total, "pos": pos, "neg": neg,
            "pos_pct": round(pos * 100 / max(total, 1)),
            "momentum": momentum, "latest": latest,
            "this_week": tw, "last_30d": _30d_count,
        }

    # ────────────────────────────────────────────────────────────────────────
    # SECTION 1: WHAT THE MARKET WANTS
    # ────────────────────────────────────────────────────────────────────────
    st.markdown("#### What the Market Wants")
    st.caption("Unmet buyer needs ranked by signal volume and source diversity.")

    # Item 15: Minimum 3-signal threshold for main display
    _strong_opps = {opp: od for opp, od in opportunity_data.items() if od["evidence"] >= 3}
    _weak_opps = {opp: od for opp, od in opportunity_data.items() if 1 <= od["evidence"] < 3}
    _sorted_strong = sorted(_strong_opps.items(), key=lambda x: x[1]["evidence"], reverse=True)

    for opp, od in _sorted_strong:
        _recent_ct = sum(1 for s in od["signals"] if s.get("post_date", "") >= _30d_ago)
        _top_sigs = sorted(od["signals"], key=lambda x: _relevance_score(x), reverse=True)[:3]

        # Build brief from top signals
        _buyer_quotes = [s for s in od["signals"] if s.get("is_buyer_voice") or s.get("is_feature_request")]
        if _buyer_quotes:
            _brief = _relevance_sentence(_buyer_quotes[0]) or "Buyers are actively seeking this capability."
        elif _top_sigs:
            _brief = _relevance_sentence(_top_sigs[0]) or "Market signals indicate demand for this capability."
        else:
            _brief = "Signal data supports this as an emerging market need."

        with st.container(border=True):
            _opp_c1, _opp_c2 = st.columns([4, 1])
            with _opp_c1:
                st.markdown(f"**{opp}**")
                st.caption(f"{od['confidence']}% confidence, {od['evidence']} signals, {_recent_ct} in last 30 days")
                st.markdown(f'<div style="font-size:13px; color:#6b7280; margin-bottom:8px;">{_brief}</div>', unsafe_allow_html=True)

                # Supporting article links (2-3 max)
                for _ts in _top_sigs[:3]:
                    _ts_title = (_ts.get("title") or _ts.get("text", ""))[:100]
                    _ts_url = _ts.get("url", "")
                    _ts_src = _source_badge(_ts.get("source", ""))
                    _ts_link = f"[{_ts_title}]({_ts_url})" if _ts_url else _ts_title
                    st.markdown(f"  `{_ts_src}` {_ts_link}")

            with _opp_c2:
                # Generate PRD button inline
                if st.button("Generate PRD", key=f"prd_{opp}", type="secondary"):
                    st.session_state["prd_gen_topic"] = opp

    # Handle PRD generation
    if st.session_state.get("prd_gen_topic"):
        _prd_topic = st.session_state.pop("prd_gen_topic")
        if _anthropic_key:
            _prd_keywords = OPPORTUNITY_THEMES.get(_prd_topic, [])
            _prd_signals = []
            for _pi in insights:
                _pt = (_pi.get("text", "") + " " + _pi.get("title", "")).lower()
                if any(kw in _pt for kw in _prd_keywords):
                    _prd_signals.append(_pi)
            _prd_signals.sort(key=lambda x: _relevance_score(x), reverse=True)

            _prd_context = ""
            for _ps in _prd_signals[:12]:
                _ps_title = (_ps.get("title", "") or _ps.get("text", ""))[:100]
                _ps_src = _ps.get("source", "")
                _ps_text = _ps.get("text", "")[:200]
                _ps_sent = _ps.get("sentiment", "neutral")
                _prd_context += f'- "{_ps_title}" ({_ps_src}, sentiment: {_ps_sent})\n  "{_ps_text}"\n'

            _prd_prompt = f"""Generate a Product Requirements Document for this GEO/AEO product opportunity based on real market signals.

Opportunity: {_prd_topic}

Signal data context:
{_prd_context}

Return ONLY valid Markdown with this exact structure:

# Product Requirements Document
**Theme:** [theme name]
**Opportunity Area:** [category]
---
## 1. Problem Statement
[2-3 sentence narrative]

**Verbatim market signals:**
- "[quote]" ([source])
- "[quote]" ([source])

---
## 2. Goals & Success Metrics
**Goals:**
- [goal]

**Success Metrics:**
- [metric with % target]

---
## 3. User Stories
1. As a [persona], I want [need] so that [outcome]

---
## 4. Requirements
### P0 (Must Have)
- **[feature]** description
### P1 (Should Have)
- **[feature]** description
### P2 (Nice to Have)
- **[feature]** description

---
## 5. Out of Scope
- item

---
## 6. Open Questions
- question

---
*Requirements grounded in GEO/AEO market signal data.*
"""
            with st.spinner("Generating PRD..."):
                try:
                    import anthropic as _anth_prd
                    _prd_client = _anth_prd.Anthropic(api_key=_anthropic_key)
                    _prd_resp = _prd_client.messages.create(
                        model="claude-haiku-4-5-20251001",
                        system="You are a product manager generating a PRD from market intelligence data. Be specific and ground every claim in the provided signals.",
                        messages=[{"role": "user", "content": _prd_prompt}],
                        max_tokens=2000,
                        temperature=0.3,
                    )
                    _prd_content = _prd_resp.content[0].text
                    st.session_state["prd_result"] = _prd_content
                    st.session_state["prd_topic"] = _prd_topic
                except Exception as _prd_err:
                    st.error(f"PRD generation failed: {_prd_err}")
        else:
            st.warning("Anthropic API key not configured.")

    # Render PRD result if available
    if st.session_state.get("prd_result"):
        with st.container(border=True):
            st.markdown(st.session_state["prd_result"])
            _prd_topic_slug = st.session_state.get("prd_topic", "opportunity").lower().replace(" ", "_")
            st.download_button(
                "Download as Markdown",
                data=st.session_state["prd_result"],
                file_name=f"prd_{_prd_topic_slug}.md",
                mime="text/markdown",
                key="prd_dl_md",
            )

    # Item 15: Weak signals collapsed section
    if _weak_opps:
        with st.expander("Weak Signals to Watch (fewer than 3 signals)"):
            for opp, od in sorted(_weak_opps.items(), key=lambda x: x[1]["evidence"], reverse=True):
                st.markdown(
                    f'<div style="padding:6px 0; border-bottom:1px solid #E8E4D9;">'
                    f'<span style="font-weight:500;">{opp}</span> '
                    f'<span style="font-family:DM Mono,monospace; font-size:11px; color:#94a3b8;">'
                    f'{od["evidence"]} signal{"s" if od["evidence"] != 1 else ""}, '
                    f'{od["confidence"]}% confidence</span></div>',
                    unsafe_allow_html=True,
                )

    st.markdown("---")

    # ────────────────────────────────────────────────────────────────────────
    # SECTION 2: COMPETITIVE WHITE SPACE
    # ────────────────────────────────────────────────────────────────────────
    st.markdown("#### Competitive White Space")
    st.caption("Features with buyer demand where no competitor has a clear solution.")

    _ws_rows = ""
    _ws_items = sorted(
        [(opp, od) for opp, od in opportunity_data.items() if od["evidence"] >= 3],
        key=lambda x: x[1]["evidence"],
        reverse=True,
    )

    for opp, od in _ws_items:
        _has_it = list(od["companies_praised"] | od["companies_tried"])[:5]
        _gap_count = sum(
            1 for c in selected_comps
            if c not in od["companies_praised"]
            and c not in od["companies_complained"]
            and c not in od["companies_tried"]
        )
        _opp_score = round(od["evidence"] * (1 + _gap_count / max(len(selected_comps), 1)) * od["confidence"] / 100)

        if _has_it:
            _coverage = ", ".join(_has_it[:3])
            if len(_has_it) > 3:
                _coverage += f" +{len(_has_it) - 3} more"
        else:
            _coverage = "No competitor has a clear solution"

        _ws_rows += (
            f'<tr style="border-bottom:1px solid #E8E4D9;">'
            f'<td style="padding:8px 12px 8px 0; font-weight:500;">{opp}</td>'
            f'<td style="padding:8px 12px; text-align:center; font-family:DM Mono,monospace; font-size:12px;">{od["evidence"]}</td>'
            f'<td style="padding:8px 12px; font-size:13px; color:#6b7280;">{_coverage}</td>'
            f'<td style="padding:8px 0 8px 12px; text-align:center; font-family:DM Mono,monospace; font-size:12px; font-weight:600;">{_opp_score}</td>'
            f'</tr>'
        )

    if _ws_rows:
        st.markdown(
            f'<table style="width:100%; font-family:DM Sans,sans-serif; font-size:14px; border-collapse:collapse;">'
            f'<tr style="font-family:DM Mono,monospace; font-size:11px; color:#94a3b8; text-transform:uppercase; border-bottom:1px solid #D1CFBA;">'
            f'<th style="text-align:left; padding:6px 12px 6px 0;">Feature</th>'
            f'<th style="text-align:center; padding:6px 12px;">Signals</th>'
            f'<th style="text-align:left; padding:6px 12px;">Competitor Coverage</th>'
            f'<th style="text-align:center; padding:6px 0 6px 12px;">Opp Score</th>'
            f'</tr>{_ws_rows}</table>',
            unsafe_allow_html=True,
        )
    else:
        st.caption("No features currently have enough data to display.")

    st.markdown("---")

    # ────────────────────────────────────────────────────────────────────────
    # SECTION 3: COMPETITOR MOMENTUM
    # ────────────────────────────────────────────────────────────────────────
    st.markdown("#### Competitor Momentum")
    st.caption("Signal activity ranked by volume. Trend based on week-over-week change.")

    _mom_sorted = sorted(selected_comps, key=lambda c: _comp_stats[c]["total"], reverse=True)
    _mom_rows = ""
    for comp in _mom_sorted:
        cs = _comp_stats[comp]
        _trend = {"Rising": "\u2191", "Falling": "\u2193", "Stable": "\u2192"}.get(cs["momentum"], "-")
        _last_sig = _time_ago(cs["latest"]) if cs["latest"] else "N/A"
        _mom_rows += (
            f'<tr style="border-bottom:1px solid #E8E4D9;">'
            f'<td style="padding:8px 12px 8px 0; font-weight:500;">{comp}</td>'
            f'<td style="padding:8px 12px; text-align:center; font-family:DM Mono,monospace; font-size:12px;">{cs["this_week"]}</td>'
            f'<td style="padding:8px 12px; text-align:center; font-family:DM Mono,monospace; font-size:12px;">{cs["last_30d"]}</td>'
            f'<td style="padding:8px 12px; text-align:center; font-size:16px;">{_trend}</td>'
            f'<td style="padding:8px 0 8px 12px; text-align:right; font-family:DM Mono,monospace; font-size:11px; color:#94a3b8;">{_last_sig}</td>'
            f'</tr>'
        )

    st.markdown(
        f'<table style="width:100%; font-family:DM Sans,sans-serif; font-size:14px; border-collapse:collapse;">'
        f'<tr style="font-family:DM Mono,monospace; font-size:11px; color:#94a3b8; text-transform:uppercase; border-bottom:1px solid #D1CFBA;">'
        f'<th style="text-align:left; padding:6px 12px 6px 0;">Competitor</th>'
        f'<th style="text-align:center; padding:6px 12px;">This Week</th>'
        f'<th style="text-align:center; padding:6px 12px;">Last 30 Days</th>'
        f'<th style="text-align:center; padding:6px 12px;">Trend</th>'
        f'<th style="text-align:right; padding:6px 0 6px 12px;">Last Signal</th>'
        f'</tr>{_mom_rows}</table>',
        unsafe_allow_html=True,
    )

    st.markdown("---")

    # ────────────────────────────────────────────────────────────────────────
    # SECTION 4: DOWNLOAD ROADMAP BRIEF
    # ────────────────────────────────────────────────────────────────────────
    st.markdown("#### Download Roadmap Brief")
    st.caption("Export the full roadmap analysis as a formatted .docx.")

    def _export_roadmap_brief():
        """Generate a Roadmap Brief .docx."""
        doc = DocxDocument()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)

        total_signals = len(insights)
        date_display = datetime.now().strftime("%B %d, %Y")
        date_str = datetime.now().strftime("%Y-%m-%d")

        doc.add_heading("GEO Pulse Roadmap Brief", level=0)
        doc.add_paragraph(f"Generated {date_display}")
        doc.add_paragraph("")

        # Section 1: What the Market Wants
        doc.add_heading("What the Market Wants", level=1)
        for opp, od in _sorted_strong:
            _rc = sum(1 for s in od["signals"] if s.get("post_date", "") >= _30d_ago)
            doc.add_heading(opp, level=2)
            doc.add_paragraph(
                f"{od['evidence']} signals, {od['confidence']}% confidence, "
                f"{_rc} in last 30 days."
            )
            _top = sorted(od["signals"], key=lambda x: _relevance_score(x), reverse=True)[:3]
            for _ts in _top:
                _tst = (_ts.get("title") or _ts.get("text", ""))[:100]
                _tsu = _ts.get("url", "")
                p = doc.add_paragraph(style="List Bullet")
                if _tsu:
                    _docx_add_hyperlink(p, _tst, _tsu, font_size=11)
                else:
                    p.add_run(_tst)

        # Section 2: Competitive White Space
        doc.add_heading("Competitive White Space", level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        for i, h in enumerate(["Feature", "Signals", "Coverage", "Opp Score"]):
            hdr[i].text = h
        for opp, od in _ws_items:
            _has = list(od["companies_praised"] | od["companies_tried"])[:3]
            _cov = ", ".join(_has) if _has else "No clear solution"
            _gc = sum(1 for c in selected_comps
                      if c not in od["companies_praised"]
                      and c not in od["companies_complained"]
                      and c not in od["companies_tried"])
            _sc = round(od["evidence"] * (1 + _gc / max(len(selected_comps), 1)) * od["confidence"] / 100)
            row = table.add_row().cells
            row[0].text = opp
            row[1].text = str(od["evidence"])
            row[2].text = _cov
            row[3].text = str(_sc)

        # Section 3: Competitor Momentum
        doc.add_heading("Competitor Momentum", level=1)
        table2 = doc.add_table(rows=1, cols=5)
        table2.style = "Table Grid"
        hdr2 = table2.rows[0].cells
        for i, h in enumerate(["Competitor", "This Week", "Last 30 Days", "Trend", "Last Signal"]):
            hdr2[i].text = h
        for comp in _mom_sorted:
            cs = _comp_stats[comp]
            row = table2.add_row().cells
            row[0].text = comp
            row[1].text = str(cs["this_week"])
            row[2].text = str(cs["last_30d"])
            row[3].text = cs["momentum"]
            row[4].text = _time_ago(cs["latest"]) if cs["latest"] else "N/A"

        _docx_source_caption(doc, total_signals, date_str)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf

    try:
        _rm_buf = _export_roadmap_brief()
        _rm_bytes = _rm_buf.getvalue()
        _rm_date = datetime.now().strftime("%Y-%m-%d")
        st.download_button(
            label="Download Roadmap Brief (.docx)",
            data=_rm_bytes,
            file_name=f"GEOPulse_RoadmapBrief_{_rm_date}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            key="export_dl_roadmap",
            type="primary",
        )
    except Exception as _rm_ex:
        st.error(f"Export failed: {_rm_ex}")




# ---------------------------------------------------------------------------
# Footer
# ---------------------------------------------------------------------------

st.divider()
st.caption("Updated every 6 hours from Reddit, Hacker News, G2, Product Hunt, trade press, and internal sources.")
